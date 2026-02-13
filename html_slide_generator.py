"""
PDF/Image template-based slide generator.
Uses Canva PDF or image template and overlays text/images using PIL.
Supports both PDF and image templates (JPG, PNG, etc.)
"""
import os
from typing import Dict, Optional
import io
from PIL import Image, ImageDraw, ImageFont, ImageFilter, ImageEnhance
import hashlib
import img2pdf
from collections import Counter
import numpy as np
from functools import lru_cache

# PPTX support for editable Canva designs
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Fix for reportlab compatibility
try:
    hashlib.md5(b'test', usedforsecurity=False)
except TypeError:
    _original_md5 = hashlib.md5
    def _patched_md5(data=None, usedforsecurity=True):
        return _original_md5(data) if data is not None else _original_md5()
    hashlib.md5 = _patched_md5

class HTMLSlideGenerator:
    # Class-level cache for rembg session (avoids reloading model on every request)
    _rembg_session = None
    
    def __init__(self):
        from config import Config
        
        # Get the directory where this script is located (works on Render too)
        script_dir = os.path.dirname(os.path.abspath(__file__))
        # On Render: script_dir = /opt/render/project/src, project_root = /opt/render/project
        # Locally: script_dir = /path/to/slauson-automation, project_root = /path/to/slauson-automation
        # Check if we're in a 'src' subdirectory (Render) or directly in project root
        if os.path.basename(script_dir) == 'src':
            project_root = os.path.dirname(script_dir)
        else:
            project_root = script_dir
        
        print(f"DEBUG: script_dir={script_dir}, project_root={project_root}")

        # Optional debug: prove which bundled fonts are present at runtime (Railway/Render shells aren't always available).
        if os.getenv("DEBUG_FONTS", "").strip().lower() in ("1", "true", "yes", "y", "on"):
            try:
                fonts_dir = os.path.join(script_dir, "assets", "fonts")
                exists = os.path.exists(fonts_dir)
                print(f"DEBUG_FONTS: assets/fonts exists={exists} path={fonts_dir}")
                if exists:
                    files = sorted(os.listdir(fonts_dir))
                    print(f"DEBUG_FONTS: assets/fonts files={files}")
                    bebas_path = os.path.join(fonts_dir, "BebasNeue-Regular.ttf")
                    print(f"DEBUG_FONTS: BebasNeue-Regular.ttf exists={os.path.exists(bebas_path)} size={os.path.getsize(bebas_path) if os.path.exists(bebas_path) else 'N/A'}")
            except Exception as e:
                print(f"DEBUG_FONTS: could not inspect bundled fonts: {e}")
        
        # Slide template path - check config first, then try default locations
        self.template_path = getattr(Config, 'SLIDE_TEMPLATE_PATH', None) if hasattr(Config, 'SLIDE_TEMPLATE_PATH') else None
        if not self.template_path or not os.path.exists(self.template_path):
            # Try default locations (project root first, then templates folder, then script dir)
            default_template_paths = [
                # Project root (where user added the file)
                os.path.join(project_root, 'SLAUSON&CO.Template.pdf'),
                os.path.join(project_root, 'SLAUSON&CO.template'),
                os.path.join(project_root, 'SLAUSON&CO.Template'),
                # Templates folder
                os.path.join(project_root, 'templates', 'SLAUSON&CO.Template.pdf'),
                os.path.join(script_dir, 'templates', 'SLAUSON&CO.Template.pdf'),
                'templates/SLAUSON&CO.Template.pdf',
                # Script directory
                os.path.join(script_dir, 'SLAUSON&CO.Template.pdf'),
                os.path.join(script_dir, 'SLAUSON&CO.template'),
                # Current working directory
                'SLAUSON&CO.Template.pdf',
                'SLAUSON&CO.template',
                # Generic template names
                os.path.join(project_root, 'templates', 'template.pdf'),
                os.path.join(script_dir, 'templates', 'template.pdf'),
                'templates/template.pdf',
                'templates/template.png',
                'templates/template.jpg',
            ]
            print(f"DEBUG: Checking {len(default_template_paths)} template paths...")
            for path in default_template_paths:
                if path and os.path.exists(path):
                    self.template_path = path
                    print(f"✓ Found template at: {self.template_path}")
                    break
                else:
                    print(f"  - Not found: {path}")
        
        # Map template path - check config first, then try default locations
        self.map_template_path = getattr(Config, 'MAP_TEMPLATE_PATH', None) if hasattr(Config, 'MAP_TEMPLATE_PATH') else None
        if not self.map_template_path or not os.path.exists(self.map_template_path):
            # Try default locations (relative to script, then project root, then current working directory)
            default_map_paths = [
                os.path.join(script_dir, 'templates', 'map_template.pdf'),
                os.path.join(project_root, 'templates', 'map_template.pdf'),
                'templates/map_template.pdf',
                os.path.join(script_dir, 'templates', 'SLAUSON&CO. (1).pdf'),
                os.path.join(project_root, 'templates', 'SLAUSON&CO. (1).pdf'),
                'templates/SLAUSON&CO. (1).pdf',
            ]
            for path in default_map_paths:
                if path and os.path.exists(path):
                    self.map_template_path = path
                    print(f"✓ Found map template at: {self.map_template_path}")
                    break

    def _load_font(self, size: int, bold: bool = False, preferred_family: Optional[str] = None):
        """
        Load a font with robust fallbacks that work on Render.
        Tries common system fonts, then DejaVu (available in most containers), then default.
        """
        def _looks_italic(font_path: str) -> bool:
            try:
                base = os.path.basename(font_path).lower()
            except Exception:
                base = str(font_path).lower()
            return ("italic" in base) or ("oblique" in base) or ("slanted" in base)

        font_candidates = []
        
        # Optional: try a preferred font family first (e.g., "Trouble")
        if preferred_family:
            try:
                fam = preferred_family.strip()
            except Exception:
                fam = str(preferred_family)
            fam_lower = fam.lower()
            
            # 1) Project-bundled fonts (recommended for consistency)
            script_dir = os.path.dirname(os.path.abspath(__file__))
            bundled_dirs = [
                os.path.join(script_dir, "assets"),
                os.path.join(script_dir, "assets", "fonts"),
                os.path.join(script_dir, "assets", "font"),
                os.path.join(script_dir, "fonts"),
            ]
            bundled_names = [
                f"{fam}.ttf",
                f"{fam}.otf",
                f"{fam} Font.ttf",
                f"{fam} Font.otf",
                f"{fam.replace(' ', '')}.ttf",
                f"{fam.replace(' ', '')}.otf",
                f"{fam}-Regular.ttf",
                f"{fam}-Regular.otf",
                f"{fam}-Bold.ttf",
                f"{fam}-Bold.otf",
            ]
            for d in bundled_dirs:
                for n in bundled_names:
                    font_candidates.append(os.path.join(d, n))
            
            # 2) Common OS font folders (local dev)
            system_dirs = [
                os.path.expanduser("~/Library/Fonts"),
                "/Library/Fonts",
                "/System/Library/Fonts",
                "/usr/share/fonts",
                "/usr/local/share/fonts",
            ]
            for d in system_dirs:
                try:
                    if not os.path.isdir(d):
                        continue
                    for fname in os.listdir(d):
                        fname_lower = fname.lower()
                        if fam_lower in fname_lower and fname_lower.endswith((".ttf", ".otf", ".ttc")):
                            if "italic" in fname_lower or "oblique" in fname_lower or "slanted" in fname_lower:
                                continue
                            font_candidates.append(os.path.join(d, fname))
                except Exception:
                    pass
            
            # 3) Linux font discovery via fontconfig (if available)
            try:
                import subprocess
                result = subprocess.run(["fc-list"], capture_output=True, text=True, timeout=1)
                if result.returncode == 0:
                    for line in result.stdout.split("\n"):
                        if fam_lower in line.lower():
                            font_path = line.split(":")[0] if ":" in line else None
                            if font_path and os.path.exists(font_path):
                                if _looks_italic(font_path):
                                    continue
                                font_candidates.insert(0, font_path)
            except Exception:
                pass

        # Preferred bundled fonts on Linux containers (Render uses these)
        if bold:
            font_candidates.append("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf")
            font_candidates.append("/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf")
        font_candidates.append("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf")
        font_candidates.append("/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf")
        # Try to find fonts via fontconfig (common on Linux).
        #
        # Important: if a preferred_family is requested, do NOT insert DejaVu/Liberation
        # ahead of the preferred family candidates (Railway/Nix often returns DejaVuSerif
        # which would override BebasNeue even when it's bundled).
        if not preferred_family:
            try:
                import subprocess
                result = subprocess.run(['fc-list'], capture_output=True, text=True, timeout=1)
                if result.returncode == 0:
                    # Look for DejaVu or Liberation in font list
                    for line in result.stdout.split('\n'):
                        if 'DejaVu' in line or 'Liberation' in line:
                            font_path = line.split(':')[0] if ':' in line else None
                            if font_path and os.path.exists(font_path):
                                if _looks_italic(font_path):
                                    continue
                                if bold and 'Bold' in font_path:
                                    font_candidates.insert(0, font_path)
                                elif not bold and 'Bold' not in font_path:
                                    font_candidates.insert(0, font_path)
            except Exception:
                pass  # Fontconfig not available, continue with hardcoded paths
        
        # macOS fonts (for local development only)
        if bold:
            font_candidates.append("/System/Library/Fonts/Helvetica-Bold.ttf")
            font_candidates.append("/System/Library/Fonts/Arial Bold.ttf")
        font_candidates.append("/System/Library/Fonts/Helvetica.ttc")
        font_candidates.append("/System/Library/Fonts/Arial.ttf")

        debug_fonts = os.getenv("DEBUG_FONTS", "").strip().lower() in ("1", "true", "yes", "y", "on")

        for path in font_candidates:
            try:
                if os.path.exists(path) or not os.path.isabs(path):
                    # Only check existence for absolute paths
                    if _looks_italic(path):
                        continue
                    font = ImageFont.truetype(path, size)
                    if debug_fonts:
                        print(f"✓ Loaded font: {path} (size={size}, bold={bold})")
                    return font
            except (OSError, IOError, Exception) as e:
                # Continue to next candidate, but log *targeted* failures when debugging.
                if debug_fonts:
                    try:
                        base = os.path.basename(path).lower()
                        # Only log likely-relevant failures to avoid spam.
                        if (
                            (preferred_family and preferred_family.lower().replace(" ", "") in base.replace(" ", ""))
                            or "bebas" in base
                            or path.startswith("/app/assets/")
                        ):
                            print(f"✗ Failed to load font: {path} (size={size}, bold={bold}) err={e}")
                    except Exception:
                        pass
                continue
        
        # Fallback: Pillow bundles DejaVu fonts in most builds; use them if available.
        # This keeps text sizing consistent even in minimal containers (Railway/Nixpacks).
        try:
            import PIL

            pil_fonts_dir = os.path.join(os.path.dirname(PIL.__file__), "fonts")
            pil_font_path = os.path.join(
                pil_fonts_dir,
                "DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf",
            )
            if os.path.exists(pil_font_path):
                font = ImageFont.truetype(pil_font_path, size)
                if os.getenv("DEBUG_FONTS", "").strip() in ("1", "true", "yes", "y", "on"):
                    print(f"✓ Loaded font (Pillow bundled): {pil_font_path} (size={size}, bold={bold})")
                return font
        except Exception:
            pass

        # Final fallback: bitmap default (NOTE: ignores `size`, will look tiny).
        print(
            "Warning: Could not load any TrueType font (custom/system/Pillow-bundled). "
            f"Falling back to PIL default bitmap font (requested size={size})."
        )
        return ImageFont.load_default()

    def _parse_name_list(self, value) -> list:
        """
        Normalize a founders/co-investors field into a list of non-empty strings.
        Accepts:
        - list[str]
        - comma-separated string
        - newline-separated string
        """
        if not value:
            return []
        if isinstance(value, list):
            items = value
        else:
            s = str(value)
            # If already newline-separated, split on newlines first.
            if "\n" in s:
                items = s.splitlines()
            else:
                items = s.split(",") if "," in s else [s]
        out = []
        for it in items:
            try:
                t = str(it).strip()
            except Exception:
                t = ""
            if t:
                out.append(t)
        return out

    def _make_white_background_transparent(self, img: Image.Image, threshold: int = 245) -> Image.Image:
        """
        Convert near-white pixels to transparent.
        Useful for icon PNGs that ship with a white box background.
        """
        try:
            rgba = img.convert("RGBA")
            arr = np.array(rgba)
            rgb = arr[:, :, :3]
            alpha = arr[:, :, 3]
            mask_white = (rgb[:, :, 0] >= threshold) & (rgb[:, :, 1] >= threshold) & (rgb[:, :, 2] >= threshold) & (alpha > 0)
            arr[mask_white, 3] = 0
            return Image.fromarray(arr, mode="RGBA")
        except Exception:
            return img.convert("RGBA")

    def _load_map_pin_icon(self, project_root: str) -> Optional[Image.Image]:
        """
        Load a custom map pin icon (PNG) if available.

        Order:
        1) MAP_PIN_ICON_PATH env var (absolute or relative)
        2) assets/pin_marker.png in this repo
        """
        candidates = []
        env_path = os.getenv("MAP_PIN_ICON_PATH")
        if env_path:
            candidates.append(env_path)
        candidates.append(os.path.join(project_root, "assets", "pin_marker.png"))

        for path in candidates:
            try:
                if not path:
                    continue
                resolved = path
                if not os.path.isabs(resolved):
                    resolved = os.path.join(project_root, resolved)
                if not os.path.exists(resolved):
                    continue
                icon = Image.open(resolved).convert("RGBA")
                icon = self._make_white_background_transparent(icon)
                # Trim transparent margins if any
                bbox = icon.split()[3].getbbox()
                if bbox:
                    icon = icon.crop(bbox)
                return icon
            except Exception:
                continue
        return None

    def _load_location_label_bg(self, project_root: str) -> Optional[Image.Image]:
        """
        Load a custom background image for the map location label.

        Order:
        1) MAP_LOCATION_LABEL_BG_PATH env var (absolute or relative)
        2) assets/location_label_bg.png in this repo
        """
        candidates = []
        env_path = os.getenv("MAP_LOCATION_LABEL_BG_PATH")
        if env_path:
            candidates.append(env_path)
        candidates.append(os.path.join(project_root, "assets", "location_label_bg.png"))

        for path in candidates:
            try:
                if not path:
                    continue
                resolved = path
                if not os.path.isabs(resolved):
                    resolved = os.path.join(project_root, resolved)
                if not os.path.exists(resolved):
                    continue
                bg = Image.open(resolved).convert("RGBA")
                bg = self._make_white_background_transparent(bg)
                # Trim transparent margins if any
                bbox = bg.split()[3].getbbox()
                if bbox:
                    bg = bg.crop(bbox)
                return bg
            except Exception:
                continue
        return None

    def _resize_pill_bg(self, bg: Image.Image, target_w: int, target_h: int) -> Image.Image:
        """
        Resize a pill-shaped background without squishing the rounded ends.

        Strategy:
        - Scale background to target height (preserve aspect ratio)
        - 3-slice horizontally: left cap + stretchable center + right cap
        """
        bg = bg.convert("RGBA")
        target_w = max(1, int(target_w))
        target_h = max(1, int(target_h))

        # First, scale to target height while preserving aspect ratio.
        if bg.size[1] != target_h:
            scale = target_h / float(bg.size[1])
            scaled_w = max(1, int(bg.size[0] * scale))
            bg = bg.resize((scaled_w, target_h), Image.Resampling.LANCZOS)

        # Choose a cap width based on height (typical pill geometry).
        # Clamp so the center slice always has positive width.
        cap_w = min(max(1, int(bg.size[1] * 0.55)), max(1, (bg.size[0] // 2) - 1))

        # If the requested width is too small, just do a normal resize.
        if target_w <= cap_w * 2 + 1:
            return bg.resize((target_w, target_h), Image.Resampling.LANCZOS)

        left = bg.crop((0, 0, cap_w, target_h))
        right = bg.crop((bg.size[0] - cap_w, 0, bg.size[0], target_h))
        center = bg.crop((cap_w, 0, bg.size[0] - cap_w, target_h))

        center_w = target_w - (cap_w * 2)
        center = center.resize((center_w, target_h), Image.Resampling.LANCZOS)

        out = Image.new("RGBA", (target_w, target_h), (0, 0, 0, 0))
        out.paste(left, (0, 0), left)
        out.paste(center, (cap_w, 0), center)
        out.paste(right, (cap_w + center_w, 0), right)
        return out

    def _dominant_visible_rgb(self, img: Image.Image) -> Optional[tuple]:
        """
        Estimate a dominant RGB color from an RGBA image, ignoring transparent pixels,
        near-white backgrounds, and very dark outline/shadow pixels.
        """
        try:
            rgba = img.convert("RGBA")
            arr = np.array(rgba)
            rgb = arr[:, :, :3].astype(np.int16)
            a = arr[:, :, 3].astype(np.int16)

            # Visible pixels only
            mask = a > 10

            # Exclude near-white (common icon background)
            mask &= ~((rgb[:, :, 0] > 245) & (rgb[:, :, 1] > 245) & (rgb[:, :, 2] > 245))

            # Exclude very dark pixels (often outline/shadow)
            luma = (0.2126 * rgb[:, :, 0] + 0.7152 * rgb[:, :, 1] + 0.0722 * rgb[:, :, 2])
            mask &= luma > 45

            coords = np.where(mask)
            if coords[0].size < 50:
                return None

            pixels = rgb[coords]
            med = np.median(pixels, axis=0)
            return (int(med[0]), int(med[1]), int(med[2]))
        except Exception:
            return None

    def _tint_rgba_to_color(self, img: Image.Image, target_rgb: tuple) -> Image.Image:
        """
        Tint an RGBA image towards a target color while preserving shading.
        Uses per-channel scaling based on the image's average visible color.
        """
        try:
            rgba = img.convert("RGBA")
            arr = np.array(rgba).astype(np.float32)
            rgb = arr[:, :, :3]
            a = arr[:, :, 3]

            mask = a > 0
            if mask.sum() == 0:
                return rgba

            base = rgb[mask].mean(axis=0)
            base = np.maximum(base, 1.0)

            target = np.array(target_rgb, dtype=np.float32)
            scale = target / base
            scale = np.clip(scale, 0.15, 6.0)

            rgb = np.clip(rgb * scale, 0, 255)
            arr[:, :, :3] = rgb
            return Image.fromarray(arr.astype(np.uint8), mode="RGBA")
        except Exception:
            return img.convert("RGBA")

    def _darken_rgba(self, img: Image.Image, factor: float = 0.88) -> Image.Image:
        """
        Darken an RGBA image by multiplying RGB channels (preserves alpha).
        factor < 1.0 makes it darker.
        """
        try:
            rgba = img.convert("RGBA")
            arr = np.array(rgba).astype(np.float32)
            arr[:, :, :3] = np.clip(arr[:, :, :3] * float(factor), 0, 255)
            return Image.fromarray(arr.astype(np.uint8), mode="RGBA")
        except Exception:
            return img.convert("RGBA")
    
    def _pdf_to_image(self, pdf_path: str) -> Image.Image:
        """Convert first page of PDF to PIL Image."""
        try:
            from pdf2image import convert_from_path
            images = convert_from_path(pdf_path, dpi=300, first_page=1, last_page=1)
            if images:
                return images[0].convert('RGBA')
        except ImportError:
            # Fallback: try PyPDF2 + extract images
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(pdf_path)
                page = reader.pages[0]
                # Try to extract images from PDF
                if '/XObject' in page['/Resources']:
                    xObject = page['/Resources']['/XObject'].get_object()
                    for obj in xObject:
                        if xObject[obj]['/Subtype'] == '/Image':
                            # Extract image data
                            data = xObject[obj].get_data()
                            img = Image.open(io.BytesIO(data))
                            return img.convert('RGBA')
                # If no images found, render PDF as image using alternative method
                raise ImportError("pdf2image required for PDF templates")
            except Exception as e:
                raise ImportError(
                    f"Could not convert PDF to image. Please install pdf2image: pip install pdf2image. "
                    f"Error: {e}"
                )
        except Exception as e:
            raise Exception(f"Failed to convert PDF to image: {e}")
    
    def _get_dominant_color(self, img: Image.Image, region: tuple, exclude_colors: list = None) -> tuple:
        """Get dominant color in a region."""
        x, y, w, h = region
        region_img = img.crop((x, y, x + w, y + h))
        pixels = list(region_img.getdata())
        
        if exclude_colors:
            pixels = [p for p in pixels if p[:3] not in exclude_colors]
        
        if not pixels:
            return (42, 42, 42)  # Default dark grey
        
        color_counts = Counter([p[:3] for p in pixels])
        return color_counts.most_common(1)[0][0]
    
    def _get_text_color_from_template(self, template: Image.Image, x: int, y: int, width: int, height: int) -> tuple:
        """Extract text color from template."""
        region = template.crop((x, y, x + width, y + height))
        pixels = list(region.getdata())
        
        bg_color = Counter([p[:3] for p in pixels]).most_common(1)[0][0]
        bg_brightness = sum(bg_color) / 3
        
        text_pixels = []
        for p in pixels:
            rgb = p[:3] if len(p) >= 3 else p
            brightness = sum(rgb) / 3
            if abs(brightness - bg_brightness) > 40:
                text_pixels.append(rgb)
        
        if text_pixels:
            return Counter(text_pixels).most_common(1)[0][0]
        
        return (255, 140, 0) if y < 200 else (255, 255, 255)
    
    def _remove_background_manual(self, img: Image.Image, tol: int = 38, feather: int = 2) -> Image.Image:
        """
        Flood-fill from the border to remove background (RGBA alpha=0),
        using multi-corner background samples + auto tolerance fallback.
        """
        try:
            import numpy as np
        except ImportError:
            print("   Warning: numpy not available for manual background removal")
            return img.convert("RGBA")

        img = img.convert("RGBA")
        arr = np.array(img)
        rgb = arr[..., :3].astype(np.int16)
        alpha = arr[..., 3].astype(np.uint8)
        H, W = rgb.shape[:2]

        corner_size = max(12, min(H, W) // 18)

        # Compute per-corner medians (better for gradients / non-uniform bg)
        tl = np.median(rgb[0:corner_size, 0:corner_size].reshape(-1, 3), axis=0)
        tr = np.median(rgb[0:corner_size, W - corner_size:W].reshape(-1, 3), axis=0)
        bl = np.median(rgb[H - corner_size:H, 0:corner_size].reshape(-1, 3), axis=0)
        br = np.median(rgb[H - corner_size:H, W - corner_size:W].reshape(-1, 3), axis=0)
        bgs = np.stack([tl, tr, bl, br], axis=0).astype(np.int16)

        def compute_close_mask(t):
            diffs = rgb[None, ...] - bgs[:, None, None, :]
            dist2 = (diffs * diffs).sum(axis=3)  # (4,H,W)
            dist = np.sqrt(dist2.min(axis=0))    # (H,W)
            return dist <= t

        tol_candidates = [tol, 45, 55, 65, 75] if tol < 45 else [tol, tol + 10, tol + 20]
        best = None

        from collections import deque

        for t in tol_candidates:
            close = compute_close_mask(t)

            bg_mask = np.zeros((H, W), dtype=bool)
            q = deque()

            def push(y, x):
                # Check bounds FIRST before accessing arrays
                if not (0 <= y < H and 0 <= x < W):
                    return
                if close[y, x] and not bg_mask[y, x]:
                    bg_mask[y, x] = True
                    q.append((y, x))

            # seed edges
            for x in range(W):
                push(0, x); push(H - 1, x)
            for y in range(H):
                push(y, 0); push(y, W - 1)

            # 8-neighborhood flood fill
            while q:
                y, x = q.popleft()
                for dy in (-1, 0, 1):
                    for dx in (-1, 0, 1):
                        if dy == 0 and dx == 0:
                            continue
                        push(y + dy, x + dx)

            removed_frac = bg_mask.mean()
            if 0.05 <= removed_frac <= 0.80:
                best = (t, bg_mask, removed_frac)
                break

            if best is None or abs(removed_frac - 0.30) < abs(best[2] - 0.30):
                best = (t, bg_mask, removed_frac)

        t, bg_mask, removed_frac = best
        print(f"   Manual BG removal: tol={t}, removed={removed_frac:.1%}")

        new_alpha = alpha.copy()
        new_alpha[bg_mask] = 0

        out = arr.copy()
        out[..., 3] = new_alpha
        out_img = Image.fromarray(out, "RGBA")

        if feather and feather > 0:
            a = out_img.split()[-1].filter(ImageFilter.GaussianBlur(radius=min(feather, 2)))
            out_img.putalpha(a)

        return out_img

    def _remove_background_gray(self, img: Image.Image, tol: int = 18, feather: int = 2) -> Image.Image:
        """
        Flood-fill background removal for GRAYSCALE headshots.
        Uses luminance-only distance with border sampling and auto-tuned tolerance.
        PROTECTS CENTER REGION to avoid removing the subject.
        """
        try:
            import numpy as np
        except ImportError:
            print("   Warning: numpy not available for gray background removal")
            return img.convert("RGBA")

        img = img.convert("RGBA")
        arr = np.array(img)
        H, W = arr.shape[:2]

        # Luminance (even if already grayscale, this is safe)
        lum = (0.299 * arr[..., 0] + 0.587 * arr[..., 1] + 0.114 * arr[..., 2]).astype(np.float32)
        alpha = arr[..., 3].astype(np.uint8)

        # Border-strip background samples in luminance (safer than corners for tight crops)
        bs = max(8, min(H, W) // 25)
        border = np.concatenate([
            lum[:bs, :].ravel(),
            lum[-bs:, :].ravel(),
            lum[:, :bs].ravel(),
            lum[:, -bs:].ravel(),
        ])
        bg = np.median(border)
        dist = np.abs(lum - bg)

        # PROTECT CENTER REGION - don't remove pixels in the center 60% of image
        # This prevents removing the person's face/body
        center_y_start = int(H * 0.2)
        center_y_end = int(H * 0.8)
        center_x_start = int(W * 0.2)
        center_x_end = int(W * 0.8)
        center_protection = np.zeros((H, W), dtype=bool)
        center_protection[center_y_start:center_y_end, center_x_start:center_x_end] = True

        from collections import deque

        def flood(t):
            close = dist <= t
            bg_mask = np.zeros((H, W), dtype=bool)
            q = deque()

            def push(y, x):
                # Check bounds FIRST before accessing arrays
                if not (0 <= y < H and 0 <= x < W):
                    return
                # CRITICAL: Don't flood into center region (protects subject)
                if center_protection[y, x]:
                    return
                if close[y, x] and not bg_mask[y, x]:
                    bg_mask[y, x] = True
                    q.append((y, x))

            # Only start flood from borders (not center)
            for x in range(W):
                push(0, x); push(H - 1, x)
            for y in range(H):
                push(y, 0); push(y, W - 1)

            while q:
                y, x = q.popleft()
                for dy in (-1, 0, 1):
                    for dx in (-1, 0, 1):
                        if dy == 0 and dx == 0:
                            continue
                        push(y + dy, x + dx)

            return bg_mask

        # More aggressive tolerance to match rembg performance (40-50% removal for headshots)
        tol_candidates = [8, 12, 15, 18, 20, 25, 30]
        best_mask = None
        best_score = None
        best_t = tol_candidates[0]

        for t in tol_candidates:
            mask = flood(t)
            removed = mask.mean()
            # Target 40-45% removal (normal for headshots - person is 50-60% of image)
            target_removal = 0.42
            score = abs(removed - target_removal)
            if best_score is None or score < best_score:
                best_score = score
                best_mask = mask
                best_t = t
            # Stop early if we get good removal (35-50% is ideal for headshots)
            if 0.35 <= removed <= 0.50:
                break

        removed_frac = best_mask.mean()
        
        # Safety check: ensure we didn't remove too much from center
        center_removed = (best_mask & center_protection).sum() / center_protection.sum()
        if center_removed > 0.05:  # If >5% of center was removed, that's bad
            print(f"   WARNING: Removed {center_removed:.1%} from center region - too aggressive!")
            # Use even more conservative mask
            for t in [3, 5, 8]:
                mask = flood(t)
                removed = mask.mean()
                center_rem = (mask & center_protection).sum() / center_protection.sum()
                if center_rem < 0.05 and 0.05 <= removed <= 0.20:
                    best_mask = mask
                    best_t = t
                    removed_frac = removed
                    print(f"   Using ultra-conservative tolerance: tol={best_t}, removed={removed_frac:.1%}, center={center_rem:.1%}")
                    break
        
        print(f"   Gray BG removal: tol={best_t}, removed={removed_frac:.1%}")
        
        # Final safety: if we removed too much overall (>60%), be more conservative
        # Note: 40-50% removal is normal for headshots (person is 50-60% of image)
        if removed_frac > 0.60:
            print(f"   WARNING: Removed too much overall ({removed_frac:.1%}), using minimal removal...")
            # Try with very low tolerance
            for t in [3, 5, 8]:
                mask = flood(t)
                removed = mask.mean()
                if 0.05 <= removed <= 0.50:  # Allow up to 50% removal (normal for headshots)
                    best_mask = mask
                    best_t = t
                    removed_frac = removed
                    print(f"   Using minimal tolerance: tol={best_t}, removed={removed_frac:.1%}")
                    break

        new_alpha = alpha.copy()
        new_alpha[best_mask] = 0

        out = arr.copy()
        out[..., 3] = new_alpha
        out_img = Image.fromarray(out, "RGBA")

        if feather and feather > 0:
            a = out_img.split()[-1].filter(ImageFilter.GaussianBlur(radius=min(feather, 2)))
            out_img.putalpha(a)

        return out_img

    def _alpha_stats(self, im: Image.Image):
        """Return (opaque_frac, transparent_frac, mean_alpha)."""
        a = np.array(im.split()[-1], dtype=np.uint8)
        opaque = (a > 200).mean()
        transparent = (a < 10).mean()
        return float(opaque), float(transparent), float(a.mean())

    def _remove_bg_rembg(self, rgba_img: Image.Image) -> Optional[Image.Image]:
        """
        Local ML segmentation using rembg (minimal, no post-processing).
        Uses cached session to avoid reloading model on every request.
        """
        try:
            import warnings
            # Suppress onnxruntime GPU warnings (Render doesn't have GPU)
            warnings.filterwarnings('ignore', category=UserWarning, module='onnxruntime')
            
            from rembg import remove, new_session
            
            # Use cached session if available, otherwise create new one
            if HTMLSlideGenerator._rembg_session is None:
                # Suppress stderr during model loading to reduce log noise
                import sys
                import os
                old_stderr = sys.stderr
                devnull = open(os.devnull, 'w')
                try:
                    sys.stderr = devnull
                    HTMLSlideGenerator._rembg_session = new_session('u2net')
                finally:
                    sys.stderr = old_stderr
                    devnull.close()
                print("   rembg session initialized (model loaded)")

            buf = io.BytesIO()
            rgba_img.save(buf, format="PNG")
            input_data = buf.getvalue()

            out = remove(input_data, session=HTMLSlideGenerator._rembg_session)

            out_img = Image.open(io.BytesIO(out)).convert("RGBA")
            out_img.load()
            return out_img
        except Exception as e:
            print(f"   rembg failed: {e}")
            return None

    def _to_grayscale_preserve_alpha(self, img: Image.Image) -> Image.Image:
        """
        Convert to grayscale while preserving alpha channel.
        """
        r, g, b, a = img.split()
        gray = img.convert("L")
        return Image.merge("RGBA", (gray, gray, gray, a))

    def _enforce_alpha_floor(self, img: Image.Image, floor: int = 50) -> Image.Image:
        """
        Ensure non-zero alpha pixels are at least 'floor' to avoid disappearing cutouts.
        """
        try:
            a = np.array(img.split()[-1], dtype=np.uint8)
            mask = a > 0
            a[mask] = np.maximum(a[mask], floor)
            a_img = Image.fromarray(a, 'L')
            r, g, b, _ = img.split()
            return Image.merge("RGBA", (r, g, b, a_img))
        except Exception:
            return img

    def _strengthen_alpha(self, img: Image.Image, thresh: int = 25, boost: float = 2.0) -> Image.Image:
        """
        Make weak alpha masks usable:
        - zero tiny alpha (background)
        - boost remaining alpha (foreground)
        """
        try:
            img = img.convert("RGBA")
            arr = np.array(img)
            a = arr[..., 3].astype(np.float32)
            a[a < thresh] = 0
            a = np.clip(a * boost, 0, 255)
            arr[..., 3] = a.astype(np.uint8)
            return Image.fromarray(arr, "RGBA")
        except Exception:
            return img


    def _remove_bg_openai(self, path: str) -> Optional[Image.Image]:
        """
        Optional AI background removal using OpenAI images.edit.
        Safe to call even if SDK or key is missing.
        """
        try:
            from openai import OpenAI
        except Exception:
            # SDK not installed
            return None

        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            # Try Config if available
            try:
                from config import Config
                api_key = getattr(Config, "OPENAI_API_KEY", None)
            except Exception:
                api_key = None
        if not api_key:
            return None

        try:
            client = OpenAI(api_key=api_key)
            # Important: pass a file handle so mimetype is recognized (not octet-stream)
            with open(path, "rb") as f:
                resp = client.images.edit(
                    model="gpt-image-1",
                    image=f,
                    prompt="Remove the background; return transparent PNG of the person.",
                    size="1024x1024",
                    output_format="png",
                )

            if not resp or not getattr(resp, "data", None):
                return None

            out_b64 = resp.data[0].b64_json
            if not out_b64:
                return None

            import base64
            out_bytes = base64.b64decode(out_b64)
            out_img = Image.open(io.BytesIO(out_bytes)).convert("RGBA")
            out_img.load()
            return out_img
        except Exception as e:
            print(f"   OpenAI bg removal failed: {e}")
            return None


    def _darken_edges(self, img: Image.Image) -> Image.Image:
        """
        Gently darkens only semi-transparent edge pixels to reduce white halo.
        Does NOT darken fully opaque pixels (avoids making the whole face dark).
        """
        try:
            import numpy as np
        except ImportError:
            return img
        arr = np.array(img)
        rgb = arr[..., :3].astype(np.float32)
        a = arr[..., 3].astype(np.float32) / 255.0

        # Only darken semi-transparent edge band (not fully opaque pixels)
        edge = (a > 0.02) & (a < 0.85)

        # Gentle darken: factor stays between 0.7..1.0 based on alpha
        factor = 0.7 + 0.3 * a
        rgb[edge] *= factor[edge][..., None]

        out = arr.copy()
        out[..., :3] = np.clip(rgb, 0, 255).astype(np.uint8)
        return Image.fromarray(out, "RGBA")

    def _fix_alpha_mask(self, img: Image.Image, a_min: int = 8) -> Image.Image:
        """
        Safe version: keeps largest blob, fills holes, hardens alpha.
        Has guards to avoid nuking the entire image if alpha is soft.
        """
        arr = np.array(img.convert("RGBA"))
        a = arr[..., 3].astype(np.uint8)
        H, W = a.shape

        # If almost nothing is non-zero alpha, don't touch it
        if (a > 0).mean() < 0.01:
            print("   _fix_alpha_mask: alpha mostly empty, skipping")
            return img

        fg = (a > a_min)

        # If fg is too small, lower threshold automatically (common for soft masks)
        if fg.mean() < 0.01:
            a_min2 = max(1, a_min // 2)
            fg = (a > a_min2)
            print(f"   _fix_alpha_mask: fg too small, lowering a_min {a_min}->{a_min2}")

        # Still nothing? skip.
        if fg.mean() < 0.005:
            print("   _fix_alpha_mask: fg still too small, skipping")
            return img

        visited = np.zeros((H, W), dtype=bool)
        best_coords = None
        best_size = 0
        from collections import deque

        for y0 in range(H):
            for x0 in range(W):
                if fg[y0, x0] and not visited[y0, x0]:
                    q = deque([(y0, x0)])
                    visited[y0, x0] = True
                    coords = []
                    while q:
                        y, x = q.popleft()
                        coords.append((y, x))
                        for dy in (-1, 0, 1):
                            for dx in (-1, 0, 1):
                                if dy == 0 and dx == 0:
                                    continue
                                ny, nx = y + dy, x + dx
                                if 0 <= ny < H and 0 <= nx < W and fg[ny, nx] and not visited[ny, nx]:
                                    visited[ny, nx] = True
                                    q.append((ny, nx))
                    if len(coords) > best_size:
                        best_size = len(coords)
                        best_coords = coords

        # If we failed to find a component (or it's tiny), don't touch the image
        if not best_coords or best_size < int(0.01 * H * W):
            print(f"   _fix_alpha_mask: best component too small ({best_size}), skipping")
            return img

        keep = np.zeros((H, W), dtype=bool)
        ys, xs = zip(*best_coords)
        keep[np.array(ys), np.array(xs)] = True

        # Fill holes: background regions not connected to border
        bg = ~keep
        hole = bg.copy()
        q = deque()

        def push(y, x):
            if 0 <= y < H and 0 <= x < W and hole[y, x]:
                hole[y, x] = False
                q.append((y, x))

        for x in range(W):
            push(0, x); push(H - 1, x)
        for y in range(H):
            push(y, 0); push(y, W - 1)

        while q:
            y, x = q.popleft()
            for dy in (-1, 0, 1):
                for dx in (-1, 0, 1):
                    if dy == 0 and dx == 0:
                        continue
                    push(y + dy, x + dx)

        filled = keep | hole

        new_a = a.copy()
        new_a[~filled] = 0

        # Mild harden (avoid killing soft hair)
        # map [0..255] to [0..255] with a small lift
        new_a = np.clip((new_a.astype(np.int16) - 5) * 255 // (255 - 5), 0, 255).astype(np.uint8)

        arr[..., 3] = new_a
        return Image.fromarray(arr, "RGBA")

    def _fallback_original_headshot(self, path: str, max_size: int = 1500) -> Optional[Image.Image]:
        """Return a safe grayscale RGBA version of the original image (no bg removal)."""
        try:
            img = Image.open(path).convert("RGBA")
            img.load()
            if max(img.size) > max_size:
                img.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)
            r, g, b, a = img.split()
            gray = img.convert("L")
            return Image.merge("RGBA", (gray, gray, gray, a))
        except Exception as e:
            print(f"Warning: fallback original headshot failed: {e}")
            return None


    def _resize_cover(self, im: Image.Image, target_w: int, target_h: int) -> Image.Image:
        """
        Scale UP or DOWN to completely fill (target_w, target_h) while preserving aspect ratio,
        then center-crop to exact size. (Like CSS background-size: cover)
        """
        im = im.convert("RGBA")
        w, h = im.size
        if w == 0 or h == 0:
            return im

        scale = max(target_w / w, target_h / h)  # cover => fill box
        new_w = max(1, int(round(w * scale)))
        new_h = max(1, int(round(h * scale)))

        im = im.resize((new_w, new_h), Image.Resampling.LANCZOS)

        # Center crop to exact target size
        left = max(0, (new_w - target_w) // 2)
        top = max(0, (new_h - target_h) // 2)
        return im.crop((left, top, left + target_w, top + target_h))

    def _refine_edges(self, img: Image.Image, erode_size: int = 3, blur_radius: float = 1.0) -> Image.Image:
        """
        Shrinks the alpha mask (erode) to remove halos, then softens the edge.
        """
        r, g, b, a = img.split()
        if erode_size > 0:
            a = a.filter(ImageFilter.MinFilter(erode_size))
        if blur_radius > 0:
            a = a.filter(ImageFilter.GaussianBlur(radius=blur_radius))
        img.putalpha(a)
        return img

    def _remove_bg_best_effort(self, path: str, use_api: bool) -> Image.Image:
        """
        Best-effort background removal:
        1) remove.bg API if available
        2) local rembg segmentation
        3) conservative grayscale flood-fill as last resort
        """
        img = Image.open(path).convert("RGBA")
        img.load()

        # Downscale before anything expensive
        max_size = 1500
        if max(img.size) > max_size:
            img.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)

        # 1) remove.bg API (if available)
        if use_api:
            try:
                from image_processor import ImageProcessor
                print(f"   Removing background via API for {path} ...")
                b = ImageProcessor.remove_background(path)
                if b:
                    api_img = Image.open(io.BytesIO(b)).convert("RGBA")
                    api_img.load()
                    o, t, ma = self._alpha_stats(api_img)
                    print(f"   API alpha stats: opaque={o:.2f}, transp={t:.2f}, meanA={ma:.0f}")
                    if o > 0.10 and t > 0.10:
                        return api_img
            except Exception as e:
                print(f"   API removal failed: {e}")

        # 2) OpenAI images.edit (optional, if key + SDK available) — now a backup
        openai_img = self._remove_bg_openai(path)
        if openai_img is not None:
            o, t, ma = self._alpha_stats(openai_img)
            print(f"   OpenAI alpha stats: opaque={o:.2f}, transp={t:.2f}, meanA={ma:.0f}")
            if o > 0.10 and t > 0.10:
                return openai_img

        # 3) local rembg
        rembg_img = self._remove_bg_rembg(img)
        if rembg_img is not None:
            o, t, ma = self._alpha_stats(rembg_img)
            print(f"   rembg alpha stats: opaque={o:.2f}, transp={t:.2f}, meanA={ma:.0f}")
            if o > 0.10 and t > 0.10:
                return rembg_img

        # 4) Last resort: conservative flood-fill for studio backgrounds (protects center/subject)
        ff = self._remove_background_gray(img, tol=8, feather=1)
        o, t, ma = self._alpha_stats(ff)
        print(f"   floodfill alpha stats: opaque={o:.2f}, transp={t:.2f}, meanA={ma:.0f}")
        return ff
    
    def _detect_orange_us_bbox(self, img: Image.Image):
        """
        Detect the orange US map outline bbox in the template.
        Restricts search to the TOP-RIGHT region to avoid the orange sidebar.
        Returns (x0, y0, w, h). Falls back if it fails.
        """
        try:
            W, H = img.size
            arr = np.array(img.convert("RGB"))
            
            # --- ROI: where the map is on your template ---
            # tune these if needed, but this matches your layout (map is top-right)
            x_start = int(W * 0.55)
            x_end   = int(W * 0.99)
            y_start = int(H * 0.02)
            y_end   = int(H * 0.45)
            
            roi = arr[y_start:y_end, x_start:x_end]
            r, g, b = roi[..., 0], roi[..., 1], roi[..., 2]
            
            # Orange outline is bright orange lines, not a solid block:
            # loosen thresholds a bit so thin lines still count
            orange_mask = (r > 160) & (g > 80) & (g < 210) & (b < 140)
            
            ys, xs = np.where(orange_mask)
            if len(xs) < 300:  # thin outlines -> fewer pixels; don't require 2000
                print(f"   Warning: Only {len(xs)} orange pixels in ROI, using fallback bbox")
                return (1250, 110, 620, 450)
            
            # bbox in ROI coords
            x0r, x1r = int(xs.min()), int(xs.max())
            y0r, y1r = int(ys.min()), int(ys.max())
            
            # Convert ROI bbox back to full-image coords
            x0 = x_start + x0r
            x1 = x_start + x1r
            y0 = y_start + y0r
            y1 = y_start + y1r
            
            pad = 10
            x0 = max(0, x0 - pad)
            y0 = max(0, y0 - pad)
            x1 = min(W - 1, x1 + pad)
            y1 = min(H - 1, y1 + pad)
            
            print(f"   Detected MAP bbox (ROI): ({x0}, {y0}, {x1-x0}, {y1-y0})")
            return (x0, y0, x1 - x0, y1 - y0)
        except Exception as e:
            print(f"   Warning: Error detecting orange bbox: {e}, using fallback")
            return (1250, 110, 620, 450)
    
    @lru_cache(maxsize=256)
    def _geocode_city(self, city: str):
        """
        Geocode city name to (lat, lon) using geopy.
        Returns (lat, lon) or None.
        Cached to avoid rate limits and improve performance.
        """
        try:
            from geopy.geocoders import Nominatim
            geolocator = Nominatim(user_agent="slide_pin_placer")
            # Bias to US
            loc = geolocator.geocode(f"{city}, USA", country_codes="us", timeout=3)
            if not loc:
                return None
            print(f"   Geocoded '{city}' to ({loc.latitude}, {loc.longitude})")
            return (loc.latitude, loc.longitude)
        except ImportError:
            print(f"   Warning: geopy not installed, cannot geocode '{city}'")
            return None
        except Exception as e:
            print(f"   Warning: Geocoding failed for '{city}': {e}")
            return None
    
    def _fallback_latlon(self, city: str):
        """
        Known-good lat/lon for major US cities + all state capitals.
        This is checked BEFORE geopy so we never get bad geocoder results for common cities.
        Uses real coordinates (no nudges) except for LA which is intentionally shifted.
        Note: AK/HI coordinates will clamp to contiguous US bounds unless handled specially.
        """
        # Parse city name - handle both "City, State" and "City State" formats
        parts = city.lower().split(',')
        city_part = parts[0].strip()
        state_part = parts[1].strip() if len(parts) > 1 else ""
        
        LATLON = {
            # Keep these as-is per your request (even though LA is intentionally shifted)
            "los angeles": (34.0522, -118.5),  # intentionally shifted
            "san francisco": (37.7749, -122.4194),
            "miami": (25.7617, -80.1918),
            "new york": (40.7128, -74.0060),
            "new york city": (40.7128, -74.0060),
            
            # Major cities
            "seattle": (47.6062, -122.3321),
            "portland": (45.5152, -122.6784),
            "boston": (40.8, -74.5),  # Nudged south and west to align pin on template map
            "chicago": (41.8781, -87.6298),
            "houston": (29.7604, -95.3698),
            "dallas": (32.7767, -96.7970),
            "san antonio": (29.4241, -98.4936),
            "san diego": (32.7157, -117.1611),
            "san jose": (37.3382, -121.8863),
            "detroit": (42.3314, -83.0458),
            "minneapolis": (44.9778, -93.2650),
            "new orleans": (29.9511, -90.0715),
            "charlotte": (35.2271, -80.8431),
            "las vegas": (36.1699, -115.1398),
            "philadelphia": (39.9526, -75.1652),
            "washington": (38.9072, -77.0369),
            "washington dc": (38.9072, -77.0369),
            "dc": (38.9072, -77.0369),
            "tampa": (27.9506, -82.4572),
            "orlando": (28.5383, -81.3792),
            "pittsburgh": (40.4406, -79.9959),
            "cleveland": (41.4993, -81.6944),
            "st louis": (38.6270, -90.1994),
            "saint louis": (38.6270, -90.1994),
            "kansas city": (39.0997, -94.5786),
            "oakland": (37.8044, -122.2712),
            "fort worth": (32.7555, -97.3308),
            "memphis": (35.1495, -90.0490),
            "baltimore": (39.2904, -76.6122),
            "milwaukee": (43.0389, -87.9065),
            "tucson": (32.2226, -110.9747),
            "el paso": (31.7619, -106.4850),
            "denver": (39.7392, -104.9903),
            
            # US State Capitals (standard coords)
            "montgomery": (32.3668, -86.3000),        # AL
            "juneau": (58.3019, -134.4197),           # AK  (needs special handling on contiguous map)
            "phoenix": (33.4484, -112.0740),          # AZ
            "little rock": (34.7465, -92.2896),       # AR
            "sacramento": (38.5816, -121.4944),       # CA
            "denver": (39.7392, -104.9903),           # CO
            "hartford": (41.7658, -72.6734),          # CT
            "dover": (39.1582, -75.5244),             # DE
            "tallahassee": (30.4383, -84.2807),       # FL
            "atlanta": (33.7490, -84.3880),           # GA
            "honolulu": (21.3069, -157.8583),         # HI  (needs special handling on contiguous map)
            "boise": (43.6150, -116.2023),            # ID
            "springfield": (39.7817, -89.6501),       # IL (capital)
            "indianapolis": (39.7684, -86.1581),      # IN
            "des moines": (41.5868, -93.6250),        # IA
            "topeka": (39.0473, -95.6752),            # KS
            "frankfort": (38.2009, -84.8733),         # KY
            "baton rouge": (30.4515, -91.1871),       # LA
            "augusta": (44.3106, -69.7795),           # ME (capital)  <-- note ambiguity with Augusta, GA
            "annapolis": (38.9784, -76.4922),         # MD
            "lansing": (42.7325, -84.5555),           # MI
            "saint paul": (44.9537, -93.0900),        # MN
            "st. paul": (44.9537, -93.0900),          # MN alias
            "jackson": (32.2988, -90.1848),           # MS
            "jefferson city": (38.5767, -92.1735),    # MO
            "helena": (46.5891, -112.0391),           # MT
            "lincoln": (40.8136, -96.7026),           # NE
            "carson city": (39.1638, -119.7674),      # NV
            "concord": (43.2081, -71.5376),           # NH
            "trenton": (40.2171, -74.7429),           # NJ
            "santa fe": (35.6870, -105.9378),         # NM
            "albany": (42.6526, -73.7562),            # NY (capital)
            "raleigh": (35.7796, -78.6382),           # NC
            "bismarck": (46.8083, -100.7837),         # ND
            "columbus": (39.9612, -82.9988),          # OH
            "oklahoma city": (35.4676, -97.5164),     # OK
            "salem": (44.9429, -123.0351),            # OR
            "harrisburg": (40.2732, -76.8867),        # PA
            "providence": (41.8240, -71.4128),        # RI
            "columbia": (34.0007, -81.0348),          # SC (capital)
            "pierre": (44.3683, -100.3510),           # SD
            "nashville": (36.1627, -86.7816),         # TN
            "austin": (30.2672, -97.7431),            # TX
            "salt lake city": (40.7608, -111.8910),   # UT
            "montpelier": (44.2601, -72.5754),        # VT
            "richmond": (37.5407, -77.4360),          # VA
            "olympia": (47.0379, -122.9007),          # WA
            "charleston": (38.3498, -81.6326),        # WV (capital) -- ambiguous with Charleston, SC city
            "madison": (43.0748, -89.3848),           # WI
            "cheyenne": (41.1400, -104.8202),         # WY
            
            # Optional disambiguation aliases (HIGHLY recommended)
            "augusta me": (44.3106, -69.7795),
            "augusta ga": (33.4735, -82.0105),
            "charleston wv": (38.3498, -81.6326),
            "charleston sc": (32.7765, -79.9311),

            # ── Alabama ──
            "birmingham": (33.5207, -86.8025),
            "huntsville": (34.7304, -86.5861),
            "mobile": (30.6954, -88.0399),
            "tuscaloosa": (33.2098, -87.5692),
            "hoover": (33.4054, -86.8114),
            "dothan": (31.2232, -85.3905),
            "auburn": (32.6099, -85.4808),
            "decatur": (34.6059, -86.9833),

            # ── Alaska ──
            "anchorage": (61.2181, -149.9003),
            "fairbanks": (64.8378, -147.7164),

            # ── Arizona ──
            "mesa": (33.4152, -111.8315),
            "scottsdale": (33.4942, -111.9261),
            "chandler": (33.3062, -111.8413),
            "gilbert": (33.3528, -111.7890),
            "glendale az": (33.5387, -112.1860),
            "tempe": (33.4255, -111.9400),
            "peoria az": (33.5806, -112.2374),
            "surprise": (33.6292, -112.3322),
            "yuma": (32.6927, -114.6277),
            "avondale": (33.4356, -112.3496),
            "flagstaff": (35.1983, -111.6513),
            "goodyear": (33.4353, -112.3585),
            "lake havasu city": (34.4839, -114.3225),
            "buckeye": (33.3703, -112.5838),
            "casa grande": (32.8795, -111.7574),

            # ── Arkansas ──
            "fort smith": (35.3859, -94.3985),
            "fayetteville": (36.0626, -94.1574),
            "springdale": (36.1867, -94.1288),
            "jonesboro": (35.8423, -90.7043),
            "rogers": (36.3320, -94.1185),
            "conway": (35.0887, -92.4421),
            "north little rock": (34.7695, -92.2671),
            "bentonville": (36.3729, -94.2088),
            "pine bluff": (34.2284, -92.0032),
            "hot springs": (34.5037, -93.0552),

            # ── California ──
            "fresno": (36.7378, -119.7871),
            "long beach": (33.7701, -118.1937),
            "bakersfield": (35.3733, -119.0187),
            "anaheim": (33.8366, -117.9143),
            "santa ana": (33.7455, -117.8677),
            "riverside": (33.9806, -117.3755),
            "stockton": (37.9577, -121.2908),
            "irvine": (33.6846, -117.8265),
            "chula vista": (32.6401, -117.0842),
            "fremont": (37.5485, -121.9886),
            "santa clarita": (34.3917, -118.5426),
            "modesto": (37.6391, -120.9969),
            "moreno valley": (33.9425, -117.2297),
            "fontana": (34.0922, -117.4350),
            "glendale ca": (34.1425, -118.2551),
            "huntington beach": (33.6595, -117.9988),
            "santa rosa": (38.4405, -122.7141),
            "oxnard": (34.1975, -119.1771),
            "oceanside": (33.1959, -117.3795),
            "rancho cucamonga": (34.1064, -117.5931),
            "ontario ca": (34.0633, -117.6509),
            "garden grove": (33.7739, -117.9414),
            "elk grove": (38.4088, -121.3716),
            "corona": (33.8753, -117.5664),
            "hayward": (37.6688, -122.0808),
            "salinas": (36.6777, -121.6555),
            "pomona": (34.0551, -117.7500),
            "roseville": (38.7521, -121.2880),
            "sunnyvale": (37.3688, -122.0363),
            "escondido": (33.1192, -117.0864),
            "torrance": (33.8358, -118.3406),
            "pasadena ca": (34.1478, -118.1445),
            "visalia": (36.3302, -119.2921),
            "santa clara": (37.3541, -121.9552),
            "concord ca": (37.9780, -122.0311),
            "thousand oaks": (34.1706, -118.8376),
            "simi valley": (34.2694, -118.7815),
            "santa maria": (34.9530, -120.4357),
            "victorville": (34.5362, -117.2928),
            "berkeley": (37.8716, -122.2727),
            "vallejo": (38.1041, -122.2566),
            "el monte": (34.0686, -118.0276),
            "downey": (33.9401, -118.1332),
            "costa mesa": (33.6412, -117.9187),
            "inglewood": (33.9617, -118.3531),
            "carlsbad": (33.1581, -117.3506),
            "san buenaventura": (34.2746, -119.2290),
            "ventura": (34.2746, -119.2290),
            "fairfield ca": (38.2494, -122.0400),
            "west covina": (34.0686, -117.9390),
            "murrieta": (33.5539, -117.2139),
            "richmond ca": (37.9358, -122.3478),
            "norwalk": (33.9022, -118.0818),
            "antioch": (38.0049, -121.8058),
            "temecula": (33.4936, -117.1484),
            "burbank": (34.1808, -118.3090),
            "daly city": (37.6879, -122.4702),
            "el cajon": (32.7948, -116.9625),
            "san mateo": (37.5630, -122.3255),
            "clovis": (36.8252, -119.7029),
            "compton": (33.8959, -118.2201),
            "jurupa valley": (33.9975, -117.4856),
            "vista": (33.2000, -117.2425),
            "south gate": (33.9547, -118.2120),
            "mission viejo": (33.6000, -117.6720),
            "vacaville": (38.3566, -121.9877),
            "carson": (33.8317, -118.2820),
            "hesperia": (34.4264, -117.3009),
            "redding": (40.5865, -122.3917),
            "santa cruz": (36.9741, -122.0308),
            "lake forest ca": (33.6469, -117.6891),
            "san leandro": (37.7249, -122.1561),
            "san marcos ca": (33.1434, -117.1661),
            "whittier": (33.9792, -118.0328),
            "hawthorne": (33.9164, -118.3526),
            "citrus heights": (38.7071, -121.2811),
            "alhambra": (34.0953, -118.1270),
            "tracy": (37.7397, -121.4252),
            "livermore": (37.6819, -121.7680),
            "buena park": (33.8675, -117.9981),
            "menifee": (33.6972, -117.1851),
            "hemet": (33.7475, -116.9719),
            "lakewood ca": (33.8536, -118.1340),
            "merced": (37.3022, -120.4830),
            "chico": (39.7285, -121.8375),
            "indio": (33.7206, -116.2156),
            "redwood city": (37.4852, -122.2364),
            "lake elsinore": (33.6681, -117.3273),
            "napa": (38.2975, -122.2869),
            "tustin": (33.7458, -117.8262),
            "bellflower": (33.8817, -118.1170),
            "mountain view": (37.3861, -122.0839),
            "chino hills": (33.9898, -117.7326),
            "baldwin park": (34.0853, -117.9609),
            "alameda": (37.7652, -122.2416),
            "upland": (34.0975, -117.6484),
            "san ramon": (37.7799, -121.9780),
            "folsom": (38.6780, -121.1761),
            "pleasanton": (37.6624, -121.8747),
            "lynwood": (33.9307, -118.2115),
            "union city ca": (37.5934, -122.0438),
            "apple valley": (34.5008, -117.1859),
            "turlock": (37.4947, -120.8466),
            "perris": (33.7825, -117.2286),
            "manteca": (37.7975, -121.2161),
            "milpitas": (37.4323, -121.8996),
            "redlands": (34.0556, -117.1825),
            "woodland": (38.6785, -121.7733),
            "lodi": (38.1302, -121.2724),
            "beaumont ca": (33.9295, -116.9770),
            "san clemente": (33.4270, -117.6120),
            "palo alto": (37.4419, -122.1430),
            "yucaipa": (34.0336, -117.0431),
            "davis": (38.5449, -121.7405),
            "camarillo": (34.2164, -119.0376),
            "walnut creek": (37.9101, -122.0652),
            "madera": (36.9613, -120.0607),
            "rancho cordova": (38.5890, -121.3027),
            "san rafael": (37.9735, -122.5311),
            "cupertino": (37.3230, -122.0322),
            "santa barbara": (34.4208, -119.6982),
            "san luis obispo": (35.2828, -120.6596),
            "palm springs": (33.8303, -116.5453),
            "palm desert": (33.7222, -116.3744),
            "monterey": (36.6002, -121.8947),
            "national city": (32.6781, -117.0992),
            "rocklin": (38.7907, -121.2358),
            "petaluma": (38.2325, -122.6367),
            "brentwood ca": (37.9317, -121.6958),

            # ── Colorado ──
            "colorado springs": (38.8339, -104.8214),
            "aurora co": (39.7294, -104.8319),
            "aurora": (39.7294, -104.8319),
            "fort collins": (40.5853, -105.0844),
            "lakewood co": (39.7047, -105.0814),
            "thornton": (39.8680, -104.9719),
            "arvada": (39.8028, -105.0875),
            "pueblo": (38.2545, -104.6091),
            "westminster co": (39.8367, -105.0372),
            "boulder": (40.0150, -105.2705),
            "centennial": (39.5807, -104.8772),
            "greeley": (40.4233, -104.7091),
            "longmont": (40.1672, -105.1019),
            "loveland": (40.3978, -105.0750),
            "broomfield": (39.9205, -105.0867),
            "castle rock": (39.3722, -104.8561),
            "commerce city": (39.8083, -104.9339),
            "parker": (39.5186, -104.7614),
            "grand junction": (39.0639, -108.5506),

            # ── Connecticut ──
            "bridgeport": (41.1865, -73.1952),
            "new haven": (41.3083, -72.9279),
            "stamford": (41.0534, -73.5387),
            "waterbury": (41.5582, -73.0515),
            "norwalk ct": (41.1177, -73.4082),
            "danbury": (41.3948, -73.4540),
            "new britain": (41.6612, -72.7795),
            "bristol ct": (41.6718, -72.9493),
            "meriden": (41.5382, -72.8068),
            "milford ct": (41.2223, -73.0565),
            "west haven": (41.2706, -72.9470),

            # ── Delaware ──
            "wilmington de": (39.7391, -75.5398),
            "wilmington": (39.7391, -75.5398),
            "newark de": (39.6837, -75.7497),

            # ── Florida ──
            "jacksonville": (30.3322, -81.6557),
            "st. petersburg": (27.7676, -82.6403),
            "st petersburg": (27.7676, -82.6403),
            "saint petersburg": (27.7676, -82.6403),
            "hialeah": (25.8576, -80.2781),
            "fort lauderdale": (26.1224, -80.1373),
            "cape coral": (26.5629, -81.9495),
            "port st. lucie": (27.2730, -80.3582),
            "port st lucie": (27.2730, -80.3582),
            "pembroke pines": (26.0128, -80.2241),
            "hollywood fl": (26.0112, -80.1495),
            "gainesville fl": (29.6516, -82.3248),
            "clearwater": (27.9659, -82.8001),
            "coral springs": (26.2712, -80.2706),
            "palm bay": (28.0345, -80.5887),
            "lakeland": (28.0395, -81.9498),
            "pompano beach": (26.2379, -80.1248),
            "west palm beach": (26.7153, -80.0534),
            "davie": (26.0765, -80.2521),
            "boca raton": (26.3587, -80.0831),
            "sunrise": (26.1339, -80.1131),
            "plantation": (26.1276, -80.2331),
            "deerfield beach": (26.3185, -80.0998),
            "largo": (27.9095, -82.7873),
            "melbourne fl": (28.0836, -80.6081),
            "palm coast": (29.5846, -81.2079),
            "deltona": (28.9005, -81.2637),
            "boynton beach": (26.5254, -80.0662),
            "lauderhill": (26.1403, -80.2134),
            "weston": (26.1004, -80.3998),
            "kissimmee": (28.2920, -81.4076),
            "homestead": (25.4687, -80.4776),
            "tamarac": (26.2129, -80.2498),
            "delray beach": (26.4615, -80.0729),
            "daytona beach": (29.2108, -81.0228),
            "north port": (27.0442, -82.2360),
            "wellington": (26.6618, -80.2684),
            "jupiter": (26.9342, -80.0942),
            "ocala": (29.1872, -82.1401),
            "port orange": (29.1383, -80.9956),
            "coconut creek": (26.2517, -80.1789),
            "sanford": (28.8003, -81.2698),
            "sarasota": (27.3364, -82.5307),
            "margate": (26.2445, -80.2065),
            "pensacola": (30.4213, -87.2169),
            "fort myers": (26.6406, -81.8723),
            "coral gables": (25.7215, -80.2684),
            "bradenton": (27.4989, -82.5748),
            "panama city": (30.1588, -85.6602),
            "naples": (26.1420, -81.7948),
            "key west": (24.5551, -81.7800),
            "winter haven": (28.0222, -81.7329),
            "winter park": (28.6000, -81.3392),

            # ── Georgia ──
            "savannah": (32.0809, -81.0912),
            "columbus ga": (32.4610, -84.9877),
            "macon": (32.8407, -83.6324),
            "athens": (33.9519, -83.3576),
            "athens ga": (33.9519, -83.3576),
            "sandy springs": (33.9304, -84.3733),
            "roswell": (34.0232, -84.3616),
            "johns creek": (34.0289, -84.1986),
            "albany ga": (31.5785, -84.1557),
            "warner robins": (32.6130, -83.5988),
            "alpharetta": (34.0754, -84.2941),
            "marietta": (33.9526, -84.5499),
            "valdosta": (30.8327, -83.2785),
            "smyrna": (33.8839, -84.5144),
            "brookhaven": (33.8651, -84.3366),
            "dunwoody": (33.9462, -84.3346),
            "peachtree city": (33.3968, -84.5957),
            "kennesaw": (34.0234, -84.6155),
            "dalton": (34.7698, -84.9702),
            "gainesville ga": (34.2979, -83.8241),
            "newnan": (33.3807, -84.7997),
            "douglasville": (33.7512, -84.7477),
            "statesboro": (32.4488, -81.7832),
            "lawrenceville": (33.9562, -83.9880),
            "duluth ga": (34.0029, -84.1446),
            "hinesville": (31.8468, -81.5962),
            "rome ga": (34.2570, -85.1647),
            "peachtree corners": (33.9701, -84.2215),
            "woodstock": (34.1015, -84.5194),
            "canton ga": (34.2368, -84.4908),

            # ── Hawaii ──
            "pearl city": (21.3972, -157.9751),
            "hilo": (19.7241, -155.0868),
            "kailua": (21.4022, -157.7394),
            "waipahu": (21.3867, -158.0092),

            # ── Idaho ──
            "nampa": (43.5407, -116.5635),
            "meridian": (43.6121, -116.3915),
            "idaho falls": (43.4917, -112.0339),
            "pocatello": (42.8713, -112.4455),
            "caldwell": (43.6629, -116.6874),
            "coeur d'alene": (47.6777, -116.7805),
            "twin falls": (42.5558, -114.4701),
            "lewiston": (46.4165, -117.0177),
            "post falls": (47.7180, -116.9516),
            "eagle": (43.6955, -116.3540),

            # ── Illinois ──
            "aurora il": (41.7606, -88.3201),
            "joliet": (41.5250, -88.0817),
            "naperville": (41.7508, -88.1535),
            "rockford": (42.2711, -89.0940),
            "elgin": (42.0354, -88.2826),
            "peoria il": (40.6936, -89.5890),
            "champaign": (40.1164, -88.2434),
            "waukegan": (42.3636, -87.8448),
            "cicero": (41.8456, -87.7539),
            "bloomington il": (40.4842, -88.9937),
            "arlington heights": (42.0884, -87.9806),
            "evanston": (42.0451, -87.6877),
            "schaumburg": (42.0334, -88.0834),
            "bolingbrook": (41.6986, -88.0684),
            "decatur il": (39.8403, -88.9548),
            "palatine": (42.1103, -88.0351),
            "skokie": (42.0324, -87.7416),
            "des plaines": (42.0334, -87.8834),
            "orland park": (41.6303, -87.8584),
            "tinley park": (41.5731, -87.7845),
            "oak lawn": (41.7106, -87.7584),
            "berwyn": (41.8506, -87.7934),
            "normal": (40.5142, -88.9907),
            "wheaton": (41.8661, -88.1070),
            "mount prospect": (42.0664, -87.9373),
            "hoffman estates": (42.0420, -88.0798),
            "oak park": (41.8850, -87.7845),
            "downers grove": (41.7959, -88.0112),

            # ── Indiana ──
            "fort wayne": (41.0793, -85.1394),
            "evansville": (37.9716, -87.5711),
            "south bend": (41.6764, -86.2520),
            "carmel": (39.9784, -86.1180),
            "fishers": (39.9568, -86.0133),
            "bloomington in": (39.1653, -86.5264),
            "hammond": (41.5833, -87.5000),
            "gary": (41.5934, -87.3464),
            "lafayette in": (40.4167, -86.8753),
            "muncie": (40.1934, -85.3864),
            "terre haute": (39.4667, -87.4139),
            "kokomo": (40.4864, -86.1336),
            "noblesville": (40.0456, -86.0086),
            "anderson in": (40.1053, -85.6803),
            "greenwood": (39.6136, -86.1067),
            "elkhart": (41.6820, -85.9767),
            "mishawaka": (41.6620, -86.1586),
            "westfield in": (40.0428, -86.1275),

            # ── Iowa ──
            "cedar rapids": (41.9779, -91.6656),
            "davenport": (41.5236, -90.5776),
            "sioux city": (42.4963, -96.4049),
            "iowa city": (41.6611, -91.5302),
            "waterloo": (42.4928, -92.3426),
            "council bluffs": (41.2619, -95.8608),
            "ames": (42.0308, -93.6319),
            "dubuque": (42.5006, -90.6646),
            "west des moines": (41.5772, -93.7113),
            "ankeny": (41.7318, -93.6001),
            "urbandale": (41.6267, -93.7122),
            "cedar falls": (42.5349, -92.4455),

            # ── Kansas ──
            "wichita": (37.6872, -97.3301),
            "overland park": (38.9822, -94.6708),
            "olathe": (38.8814, -94.8192),
            "lawrence": (38.9717, -95.2353),
            "shawnee": (39.0228, -94.7152),
            "manhattan ks": (39.1836, -96.5717),
            "lenexa": (38.9536, -94.7336),
            "salina ks": (38.8403, -97.6114),
            "hutchinson ks": (38.0608, -97.9298),
            "garden city ks": (37.9717, -100.8727),
            "derby": (37.5456, -97.2689),

            # ── Kentucky ──
            "louisville": (38.2527, -85.7585),
            "lexington": (38.0406, -84.5037),
            "bowling green": (36.9685, -86.4808),
            "owensboro": (37.7719, -87.1112),
            "covington ky": (39.0837, -84.5086),
            "richmond ky": (37.7479, -84.2947),
            "georgetown ky": (38.2098, -84.5588),
            "florence ky": (38.9990, -84.6266),
            "hopkinsville": (36.8656, -87.4886),
            "nicholasville": (37.8806, -84.5730),
            "elizabethtown": (37.6940, -85.8591),
            "henderson ky": (37.8361, -87.5900),
            "paducah": (37.0834, -88.6001),

            # ── Louisiana ──
            "shreveport": (32.5252, -93.7502),
            "lafayette la": (30.2241, -92.0198),
            "lafayette": (30.2241, -92.0198),
            "lake charles": (30.2266, -93.2174),
            "kenner": (29.9941, -90.2417),
            "bossier city": (32.5160, -93.7321),
            "monroe la": (32.5093, -92.1193),
            "alexandria la": (31.3113, -92.4451),
            "houma": (29.5958, -90.7195),

            # ── Maine ──
            "portland me": (43.6591, -70.2568),
            "lewiston me": (44.1004, -70.2148),
            "bangor": (44.8016, -68.7712),
            "south portland": (43.6415, -70.2409),
            "auburn me": (44.0979, -70.2312),

            # ── Maryland ──
            "frederick": (39.4143, -77.4105),
            "rockville": (39.0840, -77.1528),
            "gaithersburg": (39.1434, -77.2014),
            "bowie": (38.9428, -76.7302),
            "hagerstown": (39.6418, -77.7200),
            "college park": (38.9807, -76.9370),
            "salisbury md": (38.3607, -75.5994),
            "laurel md": (39.0993, -76.8483),

            # ── Massachusetts ──
            "worcester": (42.2626, -71.8023),
            "springfield ma": (42.1015, -72.5898),
            "cambridge": (42.3736, -71.1097),
            "lowell": (42.6334, -71.3162),
            "brockton": (42.0834, -71.0184),
            "new bedford": (41.6362, -70.9342),
            "quincy": (42.2529, -71.0023),
            "lynn": (42.4668, -70.9495),
            "fall river": (41.7015, -71.1550),
            "newton": (42.3370, -71.2092),
            "somerville": (42.3876, -71.0995),
            "lawrence ma": (42.7070, -71.1631),
            "haverhill": (42.7762, -71.0773),
            "waltham": (42.3765, -71.2356),
            "malden": (42.4251, -71.0662),
            "medford": (42.4184, -71.1062),
            "taunton": (41.9001, -71.0898),
            "chicopee": (42.1487, -72.6079),
            "weymouth": (42.2209, -70.9395),
            "revere": (42.4084, -71.0120),
            "peabody": (42.5279, -70.9287),
            "methuen": (42.7262, -71.1909),
            "barnstable": (41.7003, -70.2962),
            "pittsfield ma": (42.4501, -73.2526),
            "attleboro": (41.9445, -71.2856),
            "arlington ma": (42.4153, -71.1565),
            "everett ma": (42.4084, -71.0537),
            "salem ma": (42.5195, -70.8967),
            "westfield ma": (42.1251, -72.7498),
            "leominster": (42.5251, -71.7598),
            "fitchburg": (42.5834, -71.8023),
            "beverly": (42.5584, -70.8801),
            "holyoke": (42.2043, -72.6162),
            "marlborough": (42.3459, -71.5523),
            "woburn": (42.4793, -71.1523),
            "chelsea": (42.3918, -71.0328),
            "amherst": (42.3751, -72.5198),

            # ── Michigan ──
            "grand rapids": (42.9634, -85.6681),
            "warren": (42.5145, -83.0147),
            "sterling heights": (42.5803, -83.0302),
            "ann arbor": (42.2808, -83.7430),
            "flint": (43.0125, -83.6875),
            "dearborn": (42.3223, -83.1763),
            "livonia": (42.3684, -83.3527),
            "troy mi": (42.6056, -83.1499),
            "westland": (42.3242, -83.4002),
            "farmington hills": (42.4853, -83.3716),
            "kalamazoo": (42.2917, -85.5872),
            "wyoming mi": (42.9134, -85.7053),
            "southfield": (42.4734, -83.2219),
            "rochester hills": (42.6584, -83.1499),
            "taylor mi": (42.2406, -83.2697),
            "pontiac": (42.6389, -83.2910),
            "st. clair shores": (42.4970, -82.8963),
            "royal oak": (42.4895, -83.1447),
            "novi": (42.4806, -83.4755),
            "dearborn heights": (42.3370, -83.2783),
            "battle creek": (42.3212, -85.1797),
            "saginaw": (43.4195, -83.9508),
            "kentwood": (42.8695, -85.6447),
            "east lansing": (42.7370, -84.4839),
            "roseville mi": (42.4973, -82.9372),
            "portage": (42.2012, -85.5800),
            "midland mi": (43.6156, -84.2472),
            "muskegon": (43.2342, -86.2484),
            "holland mi": (42.7876, -86.1089),
            "bay city mi": (43.5945, -83.8889),

            # ── Minnesota ──
            "rochester mn": (44.0121, -92.4802),
            "duluth": (46.7867, -92.1005),
            "bloomington mn": (44.8408, -93.2983),
            "brooklyn park": (45.0941, -93.3563),
            "plymouth mn": (45.0105, -93.4555),
            "maple grove": (45.0724, -93.4558),
            "woodbury mn": (44.9239, -92.9594),
            "st. cloud": (45.5608, -94.1636),
            "st cloud": (45.5608, -94.1636),
            "eagan": (44.8041, -93.1669),
            "eden prairie": (44.8547, -93.4708),
            "coon rapids": (45.1200, -93.2875),
            "burnsville": (44.7677, -93.2777),
            "blaine": (45.1608, -93.2350),
            "lakeville mn": (44.6497, -93.2427),
            "minnetonka": (44.9211, -93.4688),
            "apple valley mn": (44.7319, -93.2177),
            "edina": (44.8897, -93.3499),
            "st. louis park": (44.9483, -93.3702),
            "mankato": (44.1636, -94.0036),
            "moorhead": (46.8738, -96.7678),
            "shakopee": (44.7974, -93.5272),
            "maplewood": (44.9530, -93.0255),
            "cottage grove": (44.8277, -92.9439),
            "richfield": (44.8833, -93.2833),
            "inver grove heights": (44.8483, -93.0427),
            "andover mn": (45.2333, -93.2914),
            "savage mn": (44.7794, -93.3361),
            "brooklyn center": (45.0761, -93.3327),
            "fridley": (45.0861, -93.2633),
            "oakdale mn": (44.9633, -92.9636),
            "champlin": (45.1889, -93.3914),
            "shoreview": (45.0794, -93.1472),
            "chanhassen": (44.8622, -93.5308),
            "prior lake": (44.7133, -93.4227),
            "rosemount": (44.7394, -93.1258),
            "chaska": (44.7894, -93.6022),
            "elk river": (45.3039, -93.5672),
            "ramsey mn": (45.2614, -93.4500),
            "hastings mn": (44.7433, -92.8511),
            "white bear lake": (45.0844, -93.0097),
            "faribault": (44.2947, -93.2688),
            "lino lakes": (45.1603, -93.0886),
            "north mankato": (44.1786, -94.0337),
            "winona": (44.0500, -91.6393),
            "austin mn": (43.6666, -92.9746),
            "owatonna": (44.0839, -93.2261),
            "stillwater mn": (45.0564, -92.8061),

            # ── Mississippi ──
            "gulfport": (30.3674, -89.0928),
            "southaven": (34.9889, -90.0126),
            "biloxi": (30.3960, -88.8853),
            "hattiesburg": (31.3271, -89.2903),
            "olive branch": (34.9618, -89.8295),
            "tupelo": (34.2576, -88.7034),
            "meridian ms": (32.3643, -88.7037),
            "pearl": (32.2746, -90.0918),
            "madison ms": (32.4618, -90.1153),
            "ocean springs": (30.4113, -88.8278),
            "starkville": (33.4504, -88.8184),
            "clinton ms": (32.3415, -90.3218),
            "horn lake": (34.9551, -90.0348),
            "ridgeland": (32.4285, -90.1323),
            "brandon": (32.2732, -89.9862),
            "columbus ms": (33.4957, -88.4273),
            "vicksburg": (32.3526, -90.8779),
            "pascagoula": (30.3658, -88.5561),

            # ── Missouri ──
            "springfield mo": (37.2090, -93.2923),
            "columbia mo": (38.9517, -92.3341),
            "independence": (39.0911, -94.4155),
            "lee's summit": (38.9108, -94.3822),
            "lees summit": (38.9108, -94.3822),
            "o'fallon mo": (38.8106, -90.6998),
            "st. joseph": (39.7675, -94.8467),
            "st joseph": (39.7675, -94.8467),
            "st. charles": (38.7881, -90.4974),
            "st charles": (38.7881, -90.4974),
            "blue springs": (39.0169, -94.2816),
            "florissant": (38.7892, -90.3226),
            "joplin": (37.0842, -94.5133),
            "chesterfield mo": (38.6631, -90.5771),
            "wildwood": (38.5828, -90.6629),
            "ballwin": (38.5953, -90.5462),
            "university city": (38.6592, -90.3098),
            "cape girardeau": (37.3059, -89.5181),
            "wentzville": (38.8114, -90.8529),

            # ── Montana ──
            "billings": (45.7833, -108.5007),
            "missoula": (46.8721, -114.0001),
            "great falls": (47.5002, -111.3008),
            "bozeman": (45.6770, -111.0429),
            "butte": (46.0038, -112.5348),
            "kalispell": (48.1920, -114.3168),

            # ── Nebraska ──
            "omaha": (41.2565, -95.9345),
            "bellevue ne": (41.1544, -95.8908),
            "grand island": (40.9264, -98.3420),
            "kearney": (40.6993, -99.0832),
            "fremont ne": (41.4333, -96.4981),
            "hastings ne": (40.5861, -98.3884),
            "norfolk ne": (42.0286, -97.4170),
            "north platte": (41.1403, -100.7601),
            "columbus ne": (41.4297, -97.3684),
            "papillion": (41.1544, -96.0422),
            "la vista": (41.1836, -96.0311),
            "scottsbluff": (41.8666, -103.6672),

            # ── Nevada ──
            "reno": (39.5296, -119.8138),
            "henderson": (36.0395, -114.9817),
            "north las vegas": (36.1989, -115.1175),
            "sparks": (39.5349, -119.7527),
            "elko": (40.8324, -115.7631),

            # ── New Hampshire ──
            "manchester nh": (42.9956, -71.4548),
            "nashua": (42.7654, -71.4676),
            "dover nh": (43.1979, -70.8737),
            "rochester nh": (43.3045, -70.9756),
            "keene": (42.9337, -72.2784),
            "portsmouth nh": (43.0718, -70.7626),
            "laconia": (43.5278, -71.4704),
            "lebanon nh": (43.6423, -72.2518),
            "claremont": (43.3770, -72.3468),
            "somersworth": (43.2618, -70.8768),

            # ── New Jersey ──
            "newark": (40.7357, -74.1724),
            "jersey city": (40.7282, -74.0776),
            "paterson": (40.9168, -74.1718),
            "elizabeth": (40.6640, -74.2107),
            "clifton": (40.8584, -74.1638),
            "trenton nj": (40.2171, -74.7429),
            "camden": (39.9259, -75.1196),
            "passaic": (40.8568, -74.1285),
            "union city nj": (40.7795, -74.0246),
            "bayonne": (40.6687, -74.1143),
            "east orange": (40.7673, -74.2049),
            "vineland": (39.4863, -75.0260),
            "new brunswick": (40.4862, -74.4518),
            "hoboken": (40.7440, -74.0324),
            "perth amboy": (40.5068, -74.2654),
            "plainfield": (40.6337, -74.4074),
            "hackensack": (40.8859, -74.0435),
            "sayreville": (40.4590, -74.3607),
            "kearny": (40.7640, -74.1454),
            "linden nj": (40.6220, -74.2446),

            # ── New Mexico ──
            "albuquerque": (35.0844, -106.6504),
            "las cruces": (32.3199, -106.7637),
            "rio rancho": (35.2334, -106.6644),
            "roswell nm": (33.3943, -104.5230),
            "farmington nm": (36.7281, -108.2187),
            "hobbs": (32.7126, -103.1361),
            "clovis nm": (34.4048, -103.2052),
            "carlsbad nm": (32.4207, -104.2288),
            "alamogordo": (32.8995, -105.9603),
            "deming": (32.2687, -107.7585),
            "gallup": (35.5281, -108.7426),
            "los lunas": (34.8062, -106.7333),
            "sunland park": (31.7960, -106.5800),

            # ── New York ──
            "buffalo": (42.8864, -78.8784),
            "rochester ny": (43.1566, -77.6088),
            "yonkers": (40.9312, -73.8987),
            "syracuse": (43.0481, -76.1474),
            "white plains": (41.0340, -73.7629),
            "mount vernon ny": (40.9126, -73.8371),
            "new rochelle": (40.9115, -73.7824),
            "schenectady": (42.8142, -73.9396),
            "utica": (43.1009, -75.2327),
            "binghamton": (42.0987, -75.9180),
            "niagara falls": (43.0962, -79.0377),
            "troy ny": (42.7284, -73.6918),
            "ithaca": (42.4440, -76.5019),
            "poughkeepsie": (41.7004, -73.9210),
            "jamestown ny": (42.0970, -79.2353),
            "elmira": (42.0898, -76.8077),
            "saratoga springs": (43.0831, -73.7846),
            "kingston ny": (41.9270, -73.9974),
            "long beach ny": (40.5884, -73.6579),
            "watertown ny": (43.9748, -75.9108),
            "glen cove": (40.8623, -73.6340),

            # ── North Carolina ──
            "durham": (35.9940, -78.8986),
            "greensboro": (36.0726, -79.7920),
            "winston-salem": (36.0999, -80.2442),
            "winston salem": (36.0999, -80.2442),
            "fayetteville nc": (35.0527, -78.8784),
            "cary": (35.7915, -78.7811),
            "wilmington nc": (34.2257, -77.9447),
            "high point": (35.9557, -80.0053),
            "concord nc": (35.4088, -80.5795),
            "asheville": (35.5951, -82.5515),
            "greenville nc": (35.6127, -77.3664),
            "gastonia": (35.2621, -81.1873),
            "jacksonville nc": (34.7541, -77.4303),
            "chapel hill": (35.9132, -79.0558),
            "huntersville": (35.4107, -80.8428),
            "apex": (35.7327, -78.8503),
            "burlington nc": (36.0957, -79.4378),
            "rocky mount": (35.9382, -77.7905),
            "kannapolis": (35.4874, -80.6217),
            "mooresville": (35.5849, -80.8101),
            "wake forest": (35.9799, -78.5097),
            "new bern": (35.1085, -77.0441),
            "hickory nc": (35.7332, -81.3413),
            "indian trail": (35.0760, -80.6692),
            "cornelius": (35.4868, -80.8601),
            "holly springs": (35.6510, -78.8336),
            "sanford nc": (35.4799, -79.1803),
            "matthews nc": (35.1168, -80.7237),
            "leland": (34.2560, -78.0447),
            "kernersville": (36.1199, -80.0737),
            "monroe nc": (34.9854, -80.5495),
            "morrisville": (35.8235, -78.8256),
            "fuquay-varina": (35.5843, -78.8000),

            # ── North Dakota ──
            "fargo": (46.8772, -96.7898),
            "grand forks": (47.9253, -97.0329),
            "minot": (48.2330, -101.2923),
            "west fargo": (46.8770, -96.9003),
            "williston": (48.1470, -103.6180),
            "dickinson": (46.8792, -102.7896),
            "mandan": (46.8267, -100.8918),

            # ── Ohio ──
            "cincinnati": (39.1031, -84.5120),
            "toledo": (41.6528, -83.5379),
            "akron": (41.0814, -81.5190),
            "dayton": (39.7589, -84.1916),
            "parma": (41.4048, -81.7229),
            "canton oh": (40.7990, -81.3784),
            "youngstown": (41.0998, -80.6495),
            "lorain": (41.4528, -82.1824),
            "hamilton oh": (39.3995, -84.5613),
            "springfield oh": (39.9242, -83.8088),
            "kettering": (39.6895, -84.1688),
            "elyria": (41.3684, -82.1076),
            "lakewood oh": (41.4820, -81.7982),
            "cuyahoga falls": (41.1340, -81.4846),
            "middletown oh": (39.5151, -84.3983),
            "euclid": (41.5931, -81.5268),
            "newark oh": (40.0581, -82.4013),
            "mansfield oh": (40.7589, -82.5154),
            "mentor": (41.6662, -81.3396),
            "beavercreek": (39.7092, -84.0633),
            "strongsville": (41.3145, -81.8357),
            "dublin oh": (40.0992, -83.1141),
            "fairfield oh": (39.3454, -84.5604),
            "findlay": (41.0442, -83.6499),
            "warren oh": (41.2378, -80.8184),
            "lancaster oh": (39.7134, -82.5993),
            "lima": (40.7428, -84.1052),
            "huber heights": (39.8439, -84.1246),
            "westerville": (40.1262, -82.9291),
            "marion oh": (40.5887, -83.1285),
            "grove city oh": (39.8812, -83.0930),
            "reynoldsburg": (39.9551, -82.8121),
            "hilliard": (40.0334, -83.1588),
            "upper arlington": (40.0267, -83.0624),
            "stow": (41.1595, -81.4404),
            "north olmsted": (41.4156, -81.9232),
            "bowling green oh": (41.3748, -83.6513),
            "north royalton": (41.3137, -81.7246),
            "massillon": (40.7967, -81.5215),
            "mason oh": (39.3601, -84.3099),
            "green oh": (40.9459, -81.4846),
            "north ridgeville": (41.3892, -82.0190),
            "gahanna": (40.0192, -82.8791),
            "brunswick oh": (41.2381, -81.8418),
            "barberton": (41.0128, -81.6051),
            "wooster": (40.8051, -81.9352),
            "westlake": (41.4553, -81.9179),
            "avon lake": (41.5053, -82.0282),
            "trotwood": (39.7973, -84.3116),
            "medina oh": (41.1387, -81.8637),
            "avon oh": (41.4517, -82.0354),
            "wadsworth": (41.0256, -81.7296),
            "austintown": (41.0867, -80.7612),
            "boardman": (41.0242, -80.6623),
            "solon": (41.3895, -81.4412),
            "twinsburg": (41.3123, -81.4404),
            "fairborn": (39.8209, -84.0194),
            "perrysburg": (41.5570, -83.6271),
            "xenia": (39.6845, -83.9296),
            "delaware oh": (40.2987, -83.0680),
            "piqua": (40.1448, -84.2424),
            "troy oh": (40.0392, -84.2033),
            "sandusky": (41.4489, -82.7079),
            "oxford oh": (39.5070, -84.7452),
            "norwood oh": (39.1556, -84.4596),
            "marysville oh": (40.2362, -83.3671),

            # ── Oklahoma ──
            "tulsa": (36.1540, -95.9928),
            "norman": (35.2226, -97.4395),
            "broken arrow": (36.0609, -95.7975),
            "edmond": (35.6528, -97.4781),
            "lawton": (34.6036, -98.3959),
            "moore": (35.3395, -97.4867),
            "midwest city": (35.4495, -97.3967),
            "enid": (36.3956, -97.8784),
            "stillwater ok": (36.1156, -97.0584),
            "muskogee": (35.7479, -95.3697),
            "bartlesville": (36.7473, -95.9808),
            "owasso": (36.2695, -95.8547),
            "shawnee ok": (35.3273, -96.9253),
            "yukon ok": (35.5067, -97.7625),
            "ardmore": (34.1743, -97.1286),
            "ponca city": (36.7070, -97.0856),
            "duncan": (34.5023, -97.9578),
            "del city": (35.4420, -97.4409),
            "bixby": (35.9420, -95.8831),
            "sapulpa": (35.9987, -96.1142),
            "altus": (34.6381, -99.3340),
            "bethany": (35.5170, -97.6323),
            "mustang": (35.3842, -97.7245),
            "sand springs": (36.1398, -96.1086),
            "claremore": (36.3126, -95.6161),
            "jenks": (35.9998, -95.9681),

            # ── Oregon ──
            "eugene": (44.0521, -123.0868),
            "gresham": (45.4983, -122.4310),
            "hillsboro": (45.5229, -122.9898),
            "bend": (44.0582, -121.3153),
            "beaverton": (45.4871, -122.8037),
            "medford": (42.3265, -122.8756),
            "springfield or": (44.0462, -123.0220),
            "corvallis": (44.5646, -123.2620),
            "albany or": (44.6368, -123.1059),
            "tigard": (45.4312, -122.7715),
            "lake oswego": (45.4207, -122.6706),
            "grants pass": (42.4390, -123.3284),
            "oregon city": (45.3573, -122.6068),
            "mcminnville": (45.2101, -123.1968),
            "redmond or": (44.2726, -121.1739),
            "tualatin": (45.3840, -122.7640),
            "west linn": (45.3654, -122.6120),
            "woodburn": (45.1438, -122.8551),
            "ashland or": (42.1946, -122.7095),
            "newberg": (45.3001, -122.9726),
            "roseburg": (43.2165, -123.3417),
            "pendleton": (45.6721, -118.7886),
            "the dalles": (45.5946, -121.1787),
            "coos bay": (43.3665, -124.2179),
            "hermiston": (45.8404, -119.2895),
            "klamath falls": (42.2249, -121.7817),
            "central point": (42.3790, -122.9023),
            "canby": (45.2629, -122.6926),
            "silverton": (45.0051, -122.7834),
            "happy valley": (45.4468, -122.5318),
            "wilsonville": (45.3001, -122.7726),
            "troutdale": (45.5393, -122.3893),
            "lebanon or": (44.5368, -122.9087),
            "dallas or": (44.9193, -123.3173),
            "cottage grove": (43.7976, -123.0589),
            "monmouth": (44.8482, -123.2340),
            "stayton": (44.8007, -122.7948),
            "keizer": (44.9901, -123.0262),
            "florence or": (43.9826, -124.0998),

            # ── Pennsylvania ──
            "allentown": (40.6084, -75.4902),
            "erie": (42.1292, -80.0851),
            "reading pa": (40.3357, -75.9269),
            "scranton": (41.4090, -75.6624),
            "bethlehem pa": (40.6259, -75.3705),
            "lancaster pa": (40.0379, -76.3055),
            "levittown": (40.1551, -74.8288),
            "york pa": (39.9626, -76.7277),
            "wilkes-barre": (41.2459, -75.8813),
            "wilkes barre": (41.2459, -75.8813),
            "chester pa": (39.8496, -75.3557),
            "williamsport": (41.2412, -77.0011),
            "easton pa": (40.6884, -75.2207),
            "lebanon pa": (40.3409, -76.4114),
            "hazleton": (40.9587, -75.9746),
            "new castle pa": (41.0034, -80.3468),
            "johnstown": (40.3268, -78.9220),
            "state college": (40.7934, -77.8600),

            # ── Rhode Island ──
            "warwick": (41.7001, -71.4162),
            "cranston": (41.7798, -71.4373),
            "pawtucket": (41.8787, -71.3826),
            "east providence": (41.8137, -71.3701),
            "woonsocket": (42.0029, -71.5148),
            "newport ri": (41.4901, -71.3128),
            "central falls": (41.8906, -71.3926),

            # ── South Carolina ──
            "north charleston": (32.8546, -79.9748),
            "greenville sc": (34.8526, -82.3940),
            "rock hill": (34.9249, -81.0251),
            "mount pleasant": (32.7941, -79.8626),
            "spartanburg": (34.9496, -81.9320),
            "summerville": (33.0185, -80.1756),
            "goose creek": (32.9810, -80.0326),
            "hilton head island": (32.2163, -80.7526),
            "florence sc": (34.1954, -79.7626),
            "myrtle beach": (33.6891, -78.8867),
            "sumter": (33.9204, -80.3415),
            "anderson sc": (34.5034, -82.6501),
            "aiken": (33.5604, -81.7196),
            "greer": (34.9385, -82.2271),
            "easley": (34.8298, -82.6015),
            "simpsonville": (34.7371, -82.2543),
            "beaufort sc": (32.4316, -80.6698),
            "hanahan": (32.9185, -80.0220),
            "mauldin": (34.7788, -82.3018),
            "greenwood sc": (34.1954, -82.1618),
            "conway sc": (33.8360, -79.0478),
            "bluffton": (32.2371, -80.8604),
            "lexington sc": (33.9815, -81.2365),
            "irmo": (34.0871, -81.1832),
            "west columbia": (33.9935, -81.0740),
            "clemson": (34.6834, -82.8374),
            "cayce": (33.9660, -81.0737),
            "orangeburg": (33.4918, -80.8565),
            "seneca sc": (34.6857, -82.9532),

            # ── South Dakota ──
            "sioux falls": (43.5460, -96.7313),
            "rapid city": (44.0805, -103.2310),
            "aberdeen sd": (45.4647, -98.4865),
            "brookings": (44.3114, -96.7984),
            "watertown sd": (44.8994, -97.1150),
            "mitchell sd": (43.7094, -98.0298),
            "yankton": (42.8711, -97.3973),
            "huron": (44.3633, -98.2148),
            "vermillion": (42.7794, -96.9292),
            "spearfish": (44.4908, -103.8593),
            "box elder": (44.1125, -103.0682),

            # ── Tennessee ──
            "knoxville": (35.9606, -83.9207),
            "chattanooga": (35.0456, -85.3097),
            "clarksville": (36.5298, -87.3595),
            "murfreesboro": (35.8456, -86.3903),
            "franklin tn": (35.9251, -86.8689),
            "johnson city": (36.3134, -82.3535),
            "jackson tn": (35.6145, -88.8139),
            "bartlett": (35.2045, -89.8740),
            "hendersonville tn": (36.3048, -86.6200),
            "kingsport": (36.5484, -82.5618),
            "collierville": (35.0420, -89.6645),
            "smyrna tn": (35.9828, -86.5186),
            "cleveland tn": (35.1595, -84.8766),
            "brentwood tn": (36.0331, -86.7828),
            "spring hill tn": (35.7512, -86.9300),
            "germantown": (35.0868, -89.8101),
            "la vergne": (36.0156, -86.5819),
            "mount juliet": (36.2001, -86.5186),
            "lebanon tn": (36.2081, -86.2911),
            "gallatin": (36.3884, -86.4467),
            "cookeville": (36.1628, -85.5016),
            "oak ridge": (36.0103, -84.2697),
            "morristown tn": (36.2140, -83.2949),
            "shelbyville": (35.4834, -86.4603),
            "tullahoma": (35.3620, -86.2094),
            "sevierville": (35.8681, -83.5618),
            "maryville tn": (35.7565, -83.9705),
            "bristol tn": (36.5951, -82.1887),
            "farragut": (35.8845, -84.1539),
            "east ridge": (35.0056, -85.2519),
            "dyersburg": (36.0345, -89.3856),

            # ── Texas ──
            "arlington tx": (32.7357, -97.1081),
            "arlington": (32.7357, -97.1081),
            "corpus christi": (27.8006, -97.3964),
            "plano": (33.0198, -96.6989),
            "laredo": (27.5036, -99.5076),
            "lubbock": (33.5779, -101.8552),
            "irving": (32.8140, -96.9489),
            "garland": (32.9126, -96.6389),
            "frisco": (33.1507, -96.8236),
            "mckinney": (33.1972, -96.6397),
            "amarillo": (35.2220, -101.8313),
            "brownsville": (25.9017, -97.4975),
            "grand prairie": (32.7460, -96.9978),
            "killeen": (31.1171, -97.7278),
            "midland tx": (31.9973, -102.0779),
            "pasadena tx": (29.6911, -95.2091),
            "beaumont tx": (30.0802, -94.1266),
            "mcallen": (26.2034, -98.2300),
            "waco": (31.5493, -97.1467),
            "denton": (33.2148, -97.1331),
            "carrollton": (32.9537, -96.8903),
            "round rock": (30.5083, -97.6789),
            "abilene": (32.4487, -99.7331),
            "odessa tx": (31.8457, -102.3676),
            "pearland": (29.5636, -95.2860),
            "richardson": (32.9483, -96.7299),
            "sugar land": (29.6197, -95.6349),
            "lewisville": (33.0462, -96.9942),
            "tyler": (32.3513, -95.3011),
            "college station": (30.6280, -96.3344),
            "allen tx": (33.1032, -96.6736),
            "league city": (29.5075, -95.0949),
            "san marcos tx": (29.8833, -97.9414),
            "edinburg": (26.3017, -98.1634),
            "mission tx": (26.2159, -98.3253),
            "pharr": (26.1948, -98.1836),
            "temple": (31.0982, -97.3428),
            "flower mound": (33.0146, -97.0969),
            "new braunfels": (29.7030, -98.1245),
            "conroe": (30.3119, -95.4561),
            "north richland hills": (32.8343, -97.2289),
            "mansfield tx": (32.5632, -97.1417),
            "bryan": (30.6744, -96.3698),
            "cedar park": (30.5052, -97.8203),
            "baytown": (29.7355, -94.9774),
            "pflugerville": (30.4394, -97.6200),
            "longview": (32.5007, -94.7405),
            "harlingen": (26.1906, -97.6961),
            "wylie": (33.0151, -96.5389),
            "burleson": (32.5421, -97.3208),
            "georgetown tx": (30.6333, -97.6781),
            "texarkana": (33.4418, -94.0477),
            "missouri city": (29.6186, -95.5377),
            "rowlett": (32.9029, -96.5639),
            "cedar hill": (32.5885, -96.9561),
            "desoto": (32.5899, -96.8570),
            "wichita falls": (33.9137, -98.4934),
            "haltom city": (32.7996, -97.2689),
            "duncanville": (32.6518, -96.9086),
            "the woodlands": (30.1658, -95.4613),
            "san angelo": (31.4638, -100.4370),
            "sherman": (33.6357, -96.6089),
            "victoria": (28.8053, -96.9850),
            "lufkin": (31.3382, -94.7291),
            "nacogdoches": (31.6035, -94.6555),
            "del rio": (29.3627, -100.8968),
            "eagle pass": (28.7091, -100.4995),
            "port arthur": (29.8849, -93.9399),
            "galveston": (29.3013, -94.7977),
            "katy": (29.7858, -95.8244),
            "rockwall": (32.9310, -96.4597),
            "coppell": (32.9546, -97.0150),
            "grapevine": (32.9343, -97.0781),
            "southlake": (32.9412, -97.1342),
            "hurst": (32.8235, -97.1706),
            "euless": (32.8371, -97.0819),
            "bedford": (32.8440, -97.1431),
            "lumberton tx": (30.2655, -94.1997),
            "friendswood": (29.5294, -95.2010),
            "humble": (29.9988, -95.2622),
            "huntsville tx": (30.7235, -95.5508),
            "cleburne": (32.3476, -97.3867),
            "weatherford tx": (32.7593, -97.7972),

            # ── Utah ──
            "west valley city": (40.6916, -112.0011),
            "provo": (40.2338, -111.6585),
            "west jordan": (40.6097, -111.9391),
            "orem": (40.2969, -111.6946),
            "sandy ut": (40.5649, -111.8590),
            "sandy": (40.5649, -111.8590),
            "ogden": (41.2230, -111.9738),
            "st. george": (37.0965, -113.5684),
            "st george": (37.0965, -113.5684),
            "layton": (41.0602, -111.9711),
            "south jordan": (40.5622, -111.9297),
            "lehi": (40.3916, -111.8508),
            "millcreek": (40.6866, -111.8755),
            "taylorsville": (40.6677, -111.9388),
            "logan ut": (41.7370, -111.8338),
            "murray ut": (40.6669, -111.8880),
            "draper": (40.5246, -111.8638),
            "bountiful": (40.8894, -111.8808),
            "riverton ut": (40.5219, -111.9391),
            "roy": (41.1616, -112.0263),
            "spanish fork": (40.1149, -111.6549),
            "pleasant grove": (40.3641, -111.7385),
            "cottonwood heights": (40.6197, -111.8102),
            "tooele": (40.5308, -112.2983),
            "springville ut": (40.1652, -111.6107),
            "eagle mountain": (40.3141, -112.0097),
            "cedar city": (37.6775, -113.0619),
            "kaysville": (41.0352, -111.9386),
            "clearfield ut": (41.1105, -112.0261),
            "herriman": (40.5144, -112.0330),
            "american fork": (40.3769, -111.7952),
            "syracuse ut": (41.0894, -112.0647),
            "saratoga springs ut": (40.3491, -111.9044),
            "magna": (40.7094, -112.1016),
            "clinton ut": (41.1397, -112.0503),
            "north ogden": (41.3072, -111.9602),
            "highland": (40.4283, -111.7952),
            "centerville ut": (40.9180, -111.8722),
            "heber city": (40.5069, -111.4127),
            "payson": (40.0444, -111.7324),
            "vineyard": (40.2969, -111.7457),
            "washington ut": (37.1305, -113.5083),
            "hurricane ut": (37.1753, -113.2900),
            "ivins": (37.1684, -113.6806),
            "santa clara ut": (37.1330, -113.6539),

            # ── Vermont ──
            "burlington vt": (44.4759, -73.2121),
            "south burlington": (44.4669, -73.1710),
            "rutland": (43.6106, -72.9726),
            "barre": (44.1970, -72.5020),
            "essex junction": (44.4901, -73.1118),
            "bennington": (42.8782, -73.1968),
            "brattleboro": (42.8509, -72.5579),
            "st. johnsbury": (44.4192, -72.0151),
            "st johnsbury": (44.4192, -72.0151),
            "middlebury": (44.0153, -73.1673),
            "winooski": (44.4914, -73.1854),
            "st. albans": (44.8110, -73.0831),

            # ── Virginia ──
            "virginia beach": (36.8529, -75.9780),
            "norfolk": (36.8508, -76.2859),
            "chesapeake": (36.7682, -76.2875),
            "newport news": (37.0871, -76.4730),
            "alexandria va": (38.8048, -77.0469),
            "alexandria": (38.8048, -77.0469),
            "hampton": (37.0299, -76.3452),
            "roanoke": (37.2710, -79.9414),
            "portsmouth va": (36.8354, -76.2983),
            "suffolk": (36.7282, -76.5836),
            "lynchburg": (37.4138, -79.1422),
            "harrisonburg": (38.4496, -78.8689),
            "leesburg": (39.1157, -77.5636),
            "charlottesville": (38.0293, -78.4767),
            "danville va": (36.5860, -79.3930),
            "manassas": (38.7509, -77.4753),
            "fredericksburg": (38.3032, -77.4605),
            "staunton": (38.1496, -79.0717),
            "waynesboro": (38.0685, -78.8895),
            "winchester": (39.1857, -78.1633),
            "radford": (37.1318, -80.5764),
            "bristol va": (36.5965, -82.1885),
            "colonial heights": (37.2440, -77.4103),
            "salem va": (37.2935, -80.0548),
            "hopewell": (37.3043, -77.2872),
            "fairfax": (38.8462, -77.3064),
            "herndon": (38.9696, -77.3861),
            "vienna va": (38.9012, -77.2653),

            # ── Washington ──
            "spokane": (47.6588, -117.4260),
            "tacoma": (47.2529, -122.4443),
            "vancouver wa": (45.6387, -122.6615),
            "bellevue wa": (47.6101, -122.2015),
            "bellevue": (47.6101, -122.2015),
            "kent wa": (47.3809, -122.2348),
            "everett wa": (47.9790, -122.2021),
            "renton": (47.4829, -122.2171),
            "spokane valley": (47.6732, -117.2394),
            "federal way": (47.3223, -122.3126),
            "yakima": (46.6021, -120.5059),
            "kirkland": (47.6815, -122.2087),
            "bellingham": (48.7519, -122.4787),
            "kennewick": (46.2112, -119.1372),
            "auburn wa": (47.3073, -122.2285),
            "pasco": (46.2396, -119.1006),
            "marysville wa": (48.0518, -122.1771),
            "lakewood wa": (47.1718, -122.5185),
            "redmond wa": (47.6740, -122.1215),
            "shoreline": (47.7557, -122.3426),
            "richland": (46.2856, -119.2845),
            "burien": (47.4704, -122.3468),
            "sammamish": (47.6163, -122.0356),
            "olympia wa": (47.0379, -122.9007),
            "lacey": (47.0343, -122.8232),
            "edmonds": (47.8107, -122.3774),
            "puyallup": (47.1854, -122.2929),
            "bremerton": (47.5673, -122.6326),
            "lynnwood": (47.8209, -122.3151),
            "bothell": (47.7623, -122.2054),
            "longview wa": (46.1382, -122.9382),
            "wenatchee": (47.4235, -120.3103),
            "walla walla": (46.0646, -118.3430),
            "mount vernon wa": (48.4213, -122.3342),
            "pullman": (46.7298, -117.1817),
            "ellensburg": (46.9965, -120.5478),
            "moses lake": (47.1301, -119.2782),
            "oak harbor": (48.2932, -122.6427),
            "centralia": (46.7162, -122.9542),
            "university place": (47.2360, -122.5485),
            "tukwila": (47.4740, -122.2610),
            "seatac": (47.4436, -122.2960),
            "issaquah": (47.5301, -122.0326),
            "covington wa": (47.3582, -122.1187),
            "tumwater": (46.9710, -122.9093),
            "woodinville": (47.7543, -122.1635),
            "maple valley": (47.3929, -122.0387),
            "mercer island": (47.5707, -122.2221),
            "des moines wa": (47.4018, -122.3243),
            "kenmore": (47.7576, -122.2437),
            "snoqualmie": (47.5287, -121.8254),
            "bonney lake": (47.1771, -122.1854),
            "camas": (45.5868, -122.3998),
            "washougal": (45.5804, -122.3534),
            "battle ground wa": (45.7807, -122.5340),
            "east wenatchee": (47.4157, -120.2937),
            "cheney": (47.4874, -117.5758),
            "airway heights": (47.6457, -117.5929),

            # ── West Virginia ──
            "huntington wv": (38.4192, -82.4452),
            "huntington": (38.4192, -82.4452),
            "morgantown": (39.6295, -79.9559),
            "parkersburg": (39.2667, -81.5615),
            "wheeling": (40.0640, -80.7209),
            "weirton": (40.4190, -80.5890),
            "martinsburg": (39.4562, -77.9639),
            "fairmont wv": (39.4851, -80.1428),
            "beckley": (37.7782, -81.1882),
            "clarksburg": (39.2806, -80.3445),
            "bluefield wv": (37.2698, -81.2223),

            # ── Wisconsin ──
            "green bay": (44.5133, -88.0133),
            "kenosha": (42.5847, -87.8212),
            "racine": (42.7261, -87.7829),
            "appleton": (44.2619, -88.4154),
            "waukesha": (43.0117, -88.2315),
            "oshkosh": (44.0247, -88.5426),
            "eau claire": (44.8113, -91.4985),
            "janesville": (42.6828, -89.0187),
            "west allis": (43.0167, -88.0070),
            "brookfield wi": (43.0606, -88.1065),
            "new berlin": (42.9764, -88.1084),
            "la crosse": (43.8014, -91.2396),
            "sheboygan": (43.7508, -87.7145),
            "wauwatosa": (43.0495, -88.0076),
            "fond du lac": (43.7750, -88.4470),
            "manitowoc": (44.0886, -87.6576),
            "west bend": (43.4253, -88.1834),
            "menomonee falls": (43.1789, -88.1176),
            "sun prairie": (43.1836, -89.2137),
            "oak creek": (42.8856, -87.8631),
            "fitchburg wi": (42.9611, -89.4237),
            "franklin wi": (42.8886, -88.0384),
            "beloit": (42.5084, -89.0318),
            "greenfield wi": (42.9614, -88.0126),
            "mequon": (43.2364, -87.9845),
            "superior": (46.7208, -92.1041),
            "stevens point": (44.5236, -89.5746),
            "caledonia wi": (42.8081, -87.8240),
            "mount pleasant wi": (42.6978, -87.8762),
            "marshfield wi": (44.6689, -90.1718),
            "pleasant prairie": (42.5534, -87.9323),
            "wisconsin rapids": (44.3836, -89.8171),

            # ── Wyoming ──
            "casper": (42.8666, -106.3131),
            "laramie": (41.3114, -105.5911),
            "gillette": (44.2911, -105.5022),
            "rock springs": (41.5875, -109.2029),
            "sheridan wy": (44.7972, -106.9561),
            "green river wy": (41.5286, -109.4662),
            "evanston wy": (41.2683, -110.9632),
            "riverton wy": (42.8549, -108.3801),
            "jackson wy": (43.4799, -110.7624),
            "cody": (44.5263, -109.0565),
            "rawlins": (41.7911, -107.2387),
            "lander": (42.8330, -108.7307),
            "torrington": (42.0625, -104.1844),
            "powell": (44.7536, -108.7573),
            "douglas wy": (42.7597, -105.3822),
            "worland": (44.0169, -107.9553),
        }
        
        # Try disambiguation key first (e.g., "augusta me", "charleston wv")
        if state_part:
            # Map common state abbreviations to full state names for disambiguation
            state_map = {
                "me": "me", "maine": "me",
                "ga": "ga", "georgia": "ga",
                "wv": "wv", "west virginia": "wv",
                "sc": "sc", "south carolina": "sc"
            }
            state_key = state_map.get(state_part, "")
            if state_key:
                disambiguation_key = f"{city_part} {state_key}"
                if disambiguation_key in LATLON:
                    return LATLON[disambiguation_key]
        
        # Fall back to city name only
        key = city_part
        return LATLON.get(key)
    
    
    def _latlon_to_map_xy(self, lat: float, lon: float, map_x: int, map_y: int, map_w: int, map_h: int):
        """
        Improved lat/lon to map coordinate conversion with calibration support.
        Uses known city positions to calibrate the projection.
        """
        # More accurate contiguous US bounds
        lon_min, lon_max = -125.0, -66.5  # Slightly wider to shift pins left
        lat_min, lat_max = 24.5, 49.5
        
        # Padding to align the projection with the template's orange US outline.
        # Tuned so Dallas lands in north-central TX and SF stays on the CA coast.
        PAD_L = 0.02  # Left padding
        PAD_R = 0.03  # Right padding (reduced so eastern cities like NY pin correctly)
        PAD_T = 0.08  # Top padding
        PAD_B = 0.01  # Bottom padding
        
        inner_x = map_x + int(map_w * PAD_L)
        inner_y = map_y + int(map_h * PAD_T)
        inner_w = int(map_w * (1.0 - PAD_L - PAD_R))
        inner_h = int(map_h * (1.0 - PAD_T - PAD_B))
        
        # Clamp to bounds
        lon = max(lon_min, min(lon_max, lon))
        lat = max(lat_min, min(lat_max, lat))
        
        # Equirectangular projection
        x_norm = (lon - lon_min) / (lon_max - lon_min)
        y_norm = 1.0 - (lat - lat_min) / (lat_max - lat_min)
        
        x = inner_x + int(x_norm * inner_w)
        y = inner_y + int(y_norm * inner_h)
        
        return x, y
    
    def calibrate_map_projection(self, template_path: str):
        """
        Helper to calibrate map projection by testing known cities.
        Run this once to find the best padding values.
        """
        template = self._pdf_to_image(template_path) if template_path.endswith('.pdf') else Image.open(template_path)
        template.load()  # Load fully before resize to avoid memory issues
        template = template.resize((1920, 1080), Image.Resampling.LANCZOS)
        
        map_x, map_y, map_w, map_h = self._detect_orange_us_bbox(template)
        
        # Test cities with known positions
        test_cities = {
            'Los Angeles': (34.0522, -118.2437),
            'New York': (40.7128, -74.0060),
            'Chicago': (41.8781, -87.6298),
            'Miami': (25.7617, -80.1918),
            'Seattle': (47.6062, -122.3321),
        }
        
        print(f"\nMap bounds: x={map_x}, y={map_y}, w={map_w}, h={map_h}")
        print("\nTesting pin positions:")
        for city, (lat, lon) in test_cities.items():
            x, y = self._latlon_to_map_xy(lat, lon, map_x, map_y, map_w, map_h)
            print(f"{city:15} -> ({x}, {y})")
        
        # Draw test pins on template
        test_img = template.copy()
        draw = ImageDraw.Draw(test_img)
        for city, (lat, lon) in test_cities.items():
            x, y = self._latlon_to_map_xy(lat, lon, map_x, map_y, map_w, map_h)
            draw.ellipse([(x-10, y-10), (x+10, y+10)], fill=(255, 0, 0))
            draw.text((x+15, y-10), city, fill=(255, 0, 0))
        
        test_img.save('map_calibration_test.png')
        print("\nSaved test image to: map_calibration_test.png")
        print("Check if pins are in correct positions and adjust PAD_L, PAD_R, PAD_T, PAD_B values")
    
    def _parse_headshots(self, headshot_path):
        """
        Accepts:
          - single path string
          - comma-separated string "a.png,b.png"
          - list/tuple of paths
        Returns list[str]
        """
        if not headshot_path:
            return []
        if isinstance(headshot_path, (list, tuple)):
            return [p for p in headshot_path if p]
        if isinstance(headshot_path, str) and "," in headshot_path:
            return [p.strip() for p in headshot_path.split(",") if p.strip()]
        return [headshot_path]
    
    def _get_city_coordinates(self, city_name: str) -> tuple:
        """
        Get normalized coordinates (0-1) for US cities on a map.
        Returns (x, y) where x is 0 (west) to 1 (east), y is 0 (north) to 1 (south).
        FALLBACK ONLY - prefer using _geocode_city + _latlon_to_map_xy
        """
        # Normalize city name
        city_lower = city_name.lower().split(',')[0].strip()
        
        # Approximate US city positions (normalized coordinates) - Fallback only
        city_positions = {
            # West Coast
            'san francisco': (0.05, 0.42),
            'los angeles': (0.10, 0.60),
            'san diego': (0.10, 0.70),
            'seattle': (0.08, 0.20),
            'portland': (0.09, 0.25),
            # East Coast
            'new york': (0.92, 0.35),
            'boston': (0.90, 0.30),
            'philadelphia': (0.90, 0.38),
            'washington': (0.88, 0.40),
            'miami': (0.93, 0.80),
            'atlanta': (0.80, 0.60),
            # Central
            'chicago': (0.65, 0.35),
            'dallas': (0.48, 0.72),
            'houston': (0.50, 0.75),
            'austin': (0.48, 0.72),
            'denver': (0.40, 0.45),
            'phoenix': (0.25, 0.65),
            'las vegas': (0.20, 0.55),
            'detroit': (0.70, 0.32),
            'minneapolis': (0.55, 0.25),
            # Other major cities
            'san jose': (0.06, 0.44),
            'oakland': (0.05, 0.43),
            'sacramento': (0.07, 0.40),
        }
        
        # Try exact match first
        if city_lower in city_positions:
            return city_positions[city_lower]
        
        # Try partial matches
        for city, coords in city_positions.items():
            if city in city_lower or city_lower in city:
                return coords
        
        # Default to center of US
        return (0.50, 0.50)

    def create_slide(self, company_data: Dict, headshot_path: str, logo_path: str, map_path: Optional[str] = None, output_format: str = "pdf") -> bytes:
        # Check if template exists, try default locations if not found
        if not self.template_path or not os.path.exists(self.template_path):
            # Get the directory where this script is located
            script_dir = os.path.dirname(os.path.abspath(__file__))
            project_root = os.path.dirname(script_dir) if os.path.basename(script_dir) in ['src', 'slauson-automation'] else script_dir
            
            # Try to find template in common locations (project root first)
            possible_paths = [
                # Project root (where user added the file)
                os.path.join(project_root, 'SLAUSON&CO.Template.pdf'),
                os.path.join(project_root, 'SLAUSON&CO.template'),
                os.path.join(project_root, 'SLAUSON&CO.Template'),
                # Templates folder
                os.path.join(project_root, 'templates', 'SLAUSON&CO.Template.pdf'),
                os.path.join(script_dir, 'templates', 'SLAUSON&CO.Template.pdf'),
                'templates/SLAUSON&CO.Template.pdf',
                # Script directory
                os.path.join(script_dir, 'SLAUSON&CO.Template.pdf'),
                os.path.join(script_dir, 'SLAUSON&CO.template'),
                # Current working directory
                'SLAUSON&CO.Template.pdf',
                'SLAUSON&CO.template',
                # Generic template names
                os.path.join(project_root, 'templates', 'template.pdf'),
                os.path.join(script_dir, 'templates', 'template.pdf'),
                'templates/template.pdf',
                'templates/template.png',
                'templates/template.jpg',
            ]
            found_path = None
            for path in possible_paths:
                if path and os.path.exists(path):
                    found_path = path
                    print(f"✓ Found template at: {found_path}")
                    break
            
            if not found_path:
                raise ValueError(
                    f"Template not found. Tried: {self.template_path}\n"
                    f"Please set SLIDE_TEMPLATE_PATH in .env or place template at project root: SLAUSON&CO.Template.pdf\n"
                    f"Script dir: {script_dir}, Project root: {project_root}\n"
                    f"Checked paths: {possible_paths[:5]}..."
                )
            self.template_path = found_path
        
        # Load template (PDF or image)
        template_ext = os.path.splitext(self.template_path)[1].lower()
        if template_ext == '.pdf':
            # Convert PDF to image
            template = self._pdf_to_image(self.template_path)
        else:
            # Load as image
            template = Image.open(self.template_path).convert('RGBA')
            template.load()  # Load fully before resize to avoid memory issues
        
        # Resize to standard slide size if needed
        if template.size != (1920, 1080):
            template = template.resize((1920, 1080), Image.Resampling.LANCZOS)
        
        # Create a copy to work with
        slide = template.copy()
        draw = ImageDraw.Draw(slide)
        width, height = slide.size
        
        # Load fonts (robust fallbacks for Render/Railway)
        name_font = self._load_font(140, bold=True)
        # Body text under Founders / Co-Investors / Background: Inter 28pt Regular
        body_font = self._load_font(28, bold=False, preferred_family="Inter_28pt-Regular")
        small_font = self._load_font(20, bold=False, preferred_family="Inter_28pt-Regular")
        sidebar_font = self._load_font(28)
        
        # Extract company data
        company_name = company_data.get('name', '').upper()
        # Hard caps to avoid overflow into other regions:
        # - Founders: max 3 names
        # - Co-investors: max 4 names
        founders_items = self._parse_name_list(company_data.get("founders", ""))[:3]
        co_investors_items = self._parse_name_list(company_data.get("co_investors", ""))[:4]

        founders_text = "\n".join(founders_items)
        co_investors_text = "\n".join(co_investors_items)
        
        background_text = company_data.get('background', company_data.get('description', ''))
        # Extract city + state from address - handle formats like
        # "123 Innovation Drive, San Francisco, CA 94105" -> "San Francisco, CA"
        full_address = (company_data.get('address') or company_data.get('location') or '').strip()
        if not full_address:
            full_address = "San Francisco, CA"
        # Try to extract city + state (or fall back to whatever we were given)
        if ',' in full_address:
            parts = [p.strip() for p in full_address.split(',')]
            # If address has 3+ parts, city is usually the second part (street, city, state/zip)
            if len(parts) >= 3:
                city = parts[1]  # e.g., "San Francisco"
                # parts[2] is typically "CA 94105" (or "CA"). Keep only the state token.
                state_token = (parts[2].split() or [""])[0].strip()
                if state_token:
                    location = f"{city}, {state_token}"
                else:
                    location = city
            else:
                # If we only have 2 parts (e.g., "San Francisco, CA"), keep both.
                location = ", ".join([p for p in parts[:2] if p])
        else:
            location = full_address
        
        investment_stage = company_data.get('investment_stage', '')
        if not investment_stage:
            round_val = company_data.get('investment_round', 'PRE-SEED')
            quarter_val = company_data.get('quarter', 'Q2')
            year_val = company_data.get('year', '2024')
            investment_stage = f"{round_val} {quarter_val}, {year_val}"
        
        # 1. Replace company name (transparent background, significantly raised, very thick bold font, orange color)
        # Align with founders position
        founders_block_y = 400  # Approximate Y position of founders yellow block
        founders_text_x = 320  # Founders text X position
        
        # Detect orange US map bounding box from template (restricted to top-right ROI)
        # Do this once and reuse for both company name overlap detection and map pin placement
        map_area_x, map_area_y, map_width, map_height = self._detect_orange_us_bbox(template)
        
        # Company name position: left aligned with founders, with a consistent top margin.
        name_x = founders_text_x - 60  # Slightly more left breathing room
        desired_title_top = 95  # Keep title visually aligned regardless of font size (smaller = lower baseline)
        
        # Use orange color for title
        # Requested: #FF9100
        name_color = (255, 145, 0)
        
        # Dynamic font sizing:
        # - Try a larger default for short names (e.g., ASTRANIS) so it feels premium.
        # - If it would overlap with the map region, shrink until it fits (existing behavior).
        name_len = len(company_name.strip())
        if name_len <= 10:
            base_font_size = 230
        elif name_len <= 14:
            base_font_size = 220
        else:
            base_font_size = 185
        name_font_size = base_font_size
        
        # Calculate text width with base font to check for overlap
        # Title font: prefer Bebas Neue if available (bundled or installed), otherwise fall back.
        # Use BebasNeue-Regular.ttf (non-italic, regular face) when available.
        test_font = self._load_font(name_font_size, bold=False, preferred_family="BebasNeue")
        text_bbox = draw.textbbox((0, 0), company_name, font=test_font)
        text_width = text_bbox[2] - text_bbox[0]
        
        # If text would overlap with map, reduce font size until it fits.
        # Keep it as large as possible (less aggressive reduction than before).
        max_allowed_width = map_area_x - name_x - 90  # Leave margin before map
        while text_width > max_allowed_width and name_font_size > 60:
            # Reduce by 5% each iteration for smoother sizing (keeps text bigger)
            name_font_size = int(name_font_size * 0.95)
            test_font = self._load_font(name_font_size, bold=False, preferred_family="BebasNeue")
            text_bbox = draw.textbbox((0, 0), company_name, font=test_font)
            text_width = text_bbox[2] - text_bbox[0]
        
        # Prefer a larger minimum size, but only if it still fits.
        desired_min_size = 100
        if name_font_size < desired_min_size:
            min_font = self._load_font(desired_min_size, bold=False, preferred_family="BebasNeue")
            min_bbox = draw.textbbox((0, 0), company_name, font=min_font)
            min_width = min_bbox[2] - min_bbox[0]
            if min_width <= max_allowed_width:
                name_font_size = desired_min_size
        
        # Hard floor so it never becomes unreadably tiny
        name_font_size = max(80, name_font_size)
        if name_font_size < base_font_size:
            print(f"   Company name too long ({len(company_name)} chars), reduced font size to {name_font_size}px to avoid map overlap")
        
        # Use very thick, bold font (matching the image style - extra thick and bold)
        name_font = self._load_font(name_font_size, bold=False, preferred_family="BebasNeue")

        # Compute Y from font bbox so the visual top stays consistent across font sizes.
        try:
            final_bbox = draw.textbbox((0, 0), company_name, font=name_font)
            final_text_width = final_bbox[2] - final_bbox[0]

            # If the title did NOT need shrinking to avoid the map, nudge it right a bit so the
            # slide feels less left-heavy (e.g., "Veridyna"). Never let it overlap the map.
            if name_font_size >= base_font_size:
                right_limit_x = map_area_x - 90 - final_text_width  # keep same map margin
                available_shift = right_limit_x - name_x
                if available_shift > 0:
                    # Move a fraction of available gap, capped, with a small minimum threshold.
                    desired_shift = min(60, int(available_shift * 0.35))
                    if desired_shift >= 12:
                        name_x = int(name_x + desired_shift)
            name_y = int(desired_title_top - final_bbox[1])
        except Exception:
            name_y = 110

        if os.getenv("DEBUG_LAYOUT", "").strip() in ("1", "true", "yes", "y", "on"):
            print(
                f"DEBUG_LAYOUT title: '{company_name}' size={name_font_size} "
                f"pos=({name_x},{name_y}) maxW={max_allowed_width} textW={text_width}"
            )
        
        # Draw company name with a very subtle stroke (outline) for contrast.
        # Keep it thin + minimal passes so it doesn't look "bold".
        stroke_width = 1
        # Use alpha for a softer outline (works on RGBA slides).
        stroke_color = (200, 80, 30, 95)
        # Minimal passes to keep it very subtle.
        offsets = [
            (-stroke_width, 0),
            (stroke_width, 0),
        ]
        for dx, dy in offsets:
            draw.text((name_x + dx, name_y + dy), company_name, fill=stroke_color, font=name_font)
        
        # Then draw the main text on top
        draw.text((name_x, name_y), company_name, fill=name_color, font=name_font)
        
        # 2. Replace logo (circular, top right) - fit within circular bounds, higher position, transparent, more circular
        try:
            logo_img = Image.open(logo_path).convert('RGBA')
            # Load image fully before operations to avoid lazy loading issues
            logo_img.load()
            logo_x, logo_y = width - 170, 10  # Raised more (was 20, now 10) to avoid map overlap
            logo_size = 130
            
            # Don't erase background - keep it transparent (no black box)
            # Just paste the logo directly with circular mask
            
            # Create perfectly circular mask first
            circle_mask = Image.new('L', (logo_size, logo_size), 0)
            circle_draw = ImageDraw.Draw(circle_mask)
            circle_draw.ellipse([(0, 0), (logo_size, logo_size)], fill=255)
            
            # Resize logo to fill the entire circle (no padding, fill the circle)
            # Center crop to square first
            logo_w, logo_h = logo_img.size
            min_dim = min(logo_w, logo_h)
            logo_img = logo_img.crop(((logo_w - min_dim) // 2, (logo_h - min_dim) // 2, 
                                     (logo_w + min_dim) // 2, (logo_h + min_dim) // 2))
            
            # Resize to fill the entire circle (logo_size x logo_size)
            logo_img = logo_img.resize((logo_size, logo_size), Image.Resampling.LANCZOS)
            
            # Apply circular mask to logo - this makes it truly circular
            logo_masked = Image.new('RGBA', (logo_size, logo_size), (0, 0, 0, 0))
            logo_masked.paste(logo_img, (0, 0))  # Paste at origin, no padding
            logo_masked.putalpha(circle_mask)  # Apply circular mask to make it round
            
            slide.paste(logo_masked, (logo_x, logo_y), logo_masked)
            draw = ImageDraw.Draw(slide)
        except Exception as e:
            print(f"Warning: Could not load logo: {e}")
        
        # 3. Updated Map Logic - Map is already in template, just add location text and adjust pin position
        try:
            # Reuse map bbox detected earlier (for company name overlap detection)
            # map_area_x, map_area_y, map_width, map_height already set above
            
            # Get lat/lon for city
            # Check known-good dictionary FIRST (before geopy) so major cities always pin correctly.
            # Geopy/Nominatim sometimes returns wrong locations (e.g., "Boston MA" → New York).
            location_lower = location.lower().strip()
            if "los angeles" in location_lower or location_lower == "la":
                latlon = (34.0522, -118.5)  # Adjusted longitude to move pin left
            else:
                latlon = self._fallback_latlon(location) or self._geocode_city(location)
                if not latlon:
                    latlon = (39.5, -98.35)  # fallback center US
            
            # Debug prints to confirm bbox detection
            print(f"   DEBUG map bbox: ({map_area_x}, {map_area_y}, {map_width}, {map_height})")
            print(f"   DEBUG location: {location}, latlon: {latlon}")
            
            # Note: AK (Juneau) and HI (Honolulu) coordinates will be clamped to contiguous US bounds
            # in _latlon_to_map_xy, so they won't appear accurately on the map. Consider skipping
            # pin placement for these states or implementing inset logic if needed.
            lat, lon = latlon
            pin_x, pin_y = self._latlon_to_map_xy(lat, lon, map_area_x, map_area_y, map_width, map_height)
            # REMOVED: pin_y -= 6  # Remove aesthetic offset as it causes inaccuracy
            
            print(f"   PIN: ({pin_x}, {pin_y})")
            
            # Resolve project root once (used for custom pin + label assets)
            script_dir = os.path.dirname(os.path.abspath(__file__))
            project_root = os.path.dirname(script_dir) if os.path.basename(script_dir) == "src" else script_dir

            # Use a custom marker icon if available (recommended).
            pin_icon = self._load_map_pin_icon(project_root)
            # Force marker + label pill to requested yellow: #FFF200
            marker_rgb = (255, 242, 0)

            if pin_icon:
                # Resize while preserving aspect ratio (anchor bottom-center on pin_x/pin_y).
                target_h = 64
                scale = target_h / float(pin_icon.size[1])
                target_w = max(1, int(pin_icon.size[0] * scale))
                pin_icon = pin_icon.resize((target_w, target_h), Image.Resampling.LANCZOS)
                # Tint the icon to the requested color for consistency.
                try:
                    pin_icon = self._tint_rgba_to_color(pin_icon, marker_rgb)
                except Exception:
                    pass
                w, h = pin_icon.size
                paste_pin_x = int(pin_x - w / 2)
                paste_pin_y = int(pin_y - h)
                slide.paste(pin_icon, (paste_pin_x, paste_pin_y), pin_icon)
                draw = ImageDraw.Draw(slide)
            else:
                # Fallback: draw the pin (yellow location pin icon - teardrop shape with circular hole)
                yellow = (255, 242, 0)  # #FFF200
                black = (0, 0, 0)
                
                # Pin dimensions - larger for visibility
                pin_width = 40  # Width of the rounded top
                pin_height = 50  # Total height including point
                point_length = 12  # Length of the pointed bottom
                hole_radius = 6  # Radius of the circular hole
                
                # Create a temporary image for the pin to draw the teardrop shape
                pin_img_size = max(pin_width, pin_height + point_length) + 20
                pin_img = Image.new('RGBA', (pin_img_size, pin_img_size), (0, 0, 0, 0))
                pin_draw = ImageDraw.Draw(pin_img)
                
                # Calculate center position in the pin image
                center_x = pin_img_size // 2
                center_y = pin_img_size // 2
                
                # Draw shadow first (slightly offset and darker) for depth
                shadow_offset = 3
                shadow_center_x = center_x + shadow_offset
                shadow_center_y = center_y + shadow_offset
                
                # Shadow: rounded top (ellipse) + pointed bottom (triangle)
                # Top ellipse for shadow
                shadow_top_y = shadow_center_y - pin_height // 2
                pin_draw.ellipse(
                    [(shadow_center_x - pin_width // 2, shadow_top_y - pin_width // 4),
                     (shadow_center_x + pin_width // 2, shadow_top_y + pin_width // 4)],
                    fill=(50, 50, 50, 120)
                )
                # Bottom triangle for shadow
                shadow_points = [
                    (shadow_center_x - pin_width // 2, shadow_top_y + pin_width // 4),
                    (shadow_center_x + pin_width // 2, shadow_top_y + pin_width // 4),
                    (shadow_center_x, shadow_center_y + point_length + shadow_offset),
                ]
                pin_draw.polygon(shadow_points, fill=(50, 50, 50, 120))
                
                # Draw main teardrop shape (yellow)
                # Top: rounded ellipse (the rounded part of the teardrop)
                top_y = center_y - pin_height // 2
                pin_draw.ellipse(
                    [(center_x - pin_width // 2, top_y - pin_width // 4),
                     (center_x + pin_width // 2, top_y + pin_width // 4)],
                    fill=yellow, outline=black, width=2
                )
                
                # Bottom: triangle connecting to point (the pointed part)
                teardrop_points = [
                    (center_x - pin_width // 2, top_y + pin_width // 4),  # Left bottom of ellipse
                    (center_x + pin_width // 2, top_y + pin_width // 4),  # Right bottom of ellipse
                    (center_x, center_y + point_length),  # Point at bottom
                ]
                pin_draw.polygon(teardrop_points, fill=yellow, outline=black, width=2)
                
                # Draw circular hole in the center (dark circle for depth)
                hole_center_x = center_x
                hole_center_y = top_y  # Position hole at top center
                pin_draw.ellipse(
                    [(hole_center_x - hole_radius, hole_center_y - hole_radius),
                     (hole_center_x + hole_radius, hole_center_y + hole_radius)],
                    fill=black
                )
                
                # Add inner highlight ring for 3D effect (lighter inner ring)
                pin_draw.ellipse(
                    [(hole_center_x - hole_radius + 2, hole_center_y - hole_radius + 2),
                     (hole_center_x + hole_radius - 2, hole_center_y + hole_radius - 2)],
                    fill=(100, 100, 100)
                )
                
                # Paste the pin onto the main slide so the bottom tip lands on (pin_x, pin_y)
                paste_pin_x = pin_x - pin_img_size // 2
                paste_pin_y = pin_y - pin_img_size // 2
                # Move image up so the bottom tip hits (pin_x, pin_y) instead of centering
                paste_pin_y -= int(pin_img_size * 0.20)  # Adjust 0.18-0.25 if needed
                slide.paste(pin_img, (paste_pin_x, paste_pin_y), pin_img)
                draw = ImageDraw.Draw(slide)
            
            # Place the city name label with yellow block (similar to other yellow blocks)
            yellow = (255, 242, 0)  # #FFF200
            black = (0, 0, 0)  # Black text on yellow background
            
            # Location label font: georgiai.ttf (bundled in assets/fonts)
            location_font = self._load_font(32, bold=False, preferred_family="georgiai")
            location_stroke_width = 1
            
            # Get text dimensions with larger font
            label_bbox = draw.textbbox((0, 0), location, font=location_font, stroke_width=location_stroke_width)
            text_width = label_bbox[2] - label_bbox[0]
            text_height = label_bbox[3] - label_bbox[1]
            
            # Yellow box dimensions (bigger padding for larger text)
            box_padding_x = 24  # Increased from 20
            box_padding_y = 12  # Increased from 10
            box_width = text_width + box_padding_x * 2
            box_height = text_height + box_padding_y * 2
            
            # Position box with more space from pin; raise significantly to avoid being too low
            box_x = pin_x + 40  # More space from pin (was 15, now 40)
            box_y = pin_y - 70  # Raised more (was -50, now -70 to move label higher)
            
            # Make sure box stays within map bounds
            if box_x + box_width > map_area_x + map_width:
                box_x = pin_x - box_width - 40  # Move to the left if it would go outside
            if box_y + box_height > map_area_y + map_height:
                box_y = pin_y - box_height - 20  # Move up if it would go outside
            
            # Draw label background. Prefer custom rounded background image if available.
            label_bg = self._load_location_label_bg(project_root)
            if label_bg:
                bg_resized = self._resize_pill_bg(label_bg, int(box_width), int(box_height))
                # Tint label background to match marker color (when available)
                bg_resized = self._tint_rgba_to_color(bg_resized, marker_rgb)
                slide.paste(bg_resized, (int(box_x), int(box_y)), bg_resized)
                draw = ImageDraw.Draw(slide)
            else:
                # Fallback: plain rectangle
                draw.rectangle([(box_x, box_y), (box_x + box_width, box_y + box_height)], fill=yellow)
            
            # Draw location text in the yellow box (centered with padding)
            text_x = box_x + box_padding_x
            text_y = box_y + box_padding_y
            draw.text(
                (text_x, text_y),
                location,
                fill=black,
                font=location_font,
                stroke_width=location_stroke_width,
                stroke_fill=black
            )
            
        except Exception as e:
            print(f"Warning: Map update failed: {e}")
        
        # 4. Replace Founders text (moved left and down, transparent)
        # Detect yellow block area for founders (left side, below company name)
        founders_block_y = 400  # Approximate Y position of founders yellow block
        founders_text_y = founders_block_y + 15  # Moved down a little
        founders_text_x = 320  # Moved a little to the left (was 350)
        founders_color = self._get_text_color_from_template(template, founders_text_x, founders_text_y, 600, 100)
        
        # Don't erase background - keep it transparent (no black box)
        # Draw founders text (max 3 names)
        founders_lines = founders_text.split('\n') if founders_text else []
        for i, line in enumerate(founders_lines):
            draw.text((founders_text_x, founders_text_y + i * 35), line, fill=founders_color, font=body_font)
        
        # 5. Replace Co-Investors text (separate text box, to the right of founders, aligned with founders height)
        investors_block_y = 500  # Approximate Y position of co-investors yellow block
        investors_text_y = founders_text_y  # Aligned with founders (same height)
        investors_text_x = 650  # Moved to the right of founders (founders is at 320, so co-investors at 650)
        investors_color = self._get_text_color_from_template(template, investors_text_x, investors_text_y, 600, 100)
        
        # Don't erase background - keep it transparent (no black box)
        # Draw co-investors text (max 4 names)
        investors_lines = co_investors_text.split('\n') if co_investors_text else []
        for i, line in enumerate(investors_lines):
            draw.text((investors_text_x, investors_text_y + i * 35), line, fill=investors_color, font=body_font)
        
        # 6. Replace Background text (aligned with founders, wider text area, bigger font, transparent background)
        bg_block_y = 600  # Approximate Y position of background yellow block
        bg_text_y = bg_block_y + 5  # Text starts below yellow block (smaller gap)
        bg_text_x = founders_text_x  # Aligned with founders (moved left from 400 to match founders at 320)
        # Allow background copy to use more horizontal room (without running into the headshot on the right).
        bg_text_width = 800

        # Background font: +2px vs Founders/Co-Investors (slightly bigger)
        bg_font = self._load_font(30, bold=False, preferred_family="Inter_28pt-Regular")
        bg_line_spacing = 37

        words = background_text.split()
        lines = []
        current_line = []
        for word in words:
            test_line = ' '.join(current_line + [word])
            bbox = draw.textbbox((0, 0), test_line, font=bg_font)
            if bbox[2] - bbox[0] < bg_text_width:  # Wider width for background
                current_line.append(word)
            else:
                if current_line:
                    lines.append(' '.join(current_line))
                current_line = [word]
        if current_line:
            lines.append(' '.join(current_line))
        
        # Get text color from template (don't erase background - keep it transparent)
        bg_color = self._get_text_color_from_template(template, bg_text_x, bg_text_y, bg_text_width, 200)
        
        # Draw text directly without erasing background (transparent)
        for i, line in enumerate(lines[:10]):
            draw.text((bg_text_x, bg_text_y + i * bg_line_spacing), line, fill=bg_color, font=bg_font)
        
        # 7. Replace headshots (below map, moved left, bigger size, transparent)
        headshot_processed = False
        try:
            from image_processor import ImageProcessor
            
            # First, validate that the headshot file exists and is a valid image
            # If not, skip headshot processing (don't crash the entire slide generation)
            if not headshot_path or not os.path.exists(headshot_path):
                print(f"Warning: Headshot file not found: {headshot_path}, skipping headshot")
                headshot_path = None
            
            if headshot_path:
                # Try to open the original image first to validate it (quick check)
                try:
                    test_img = Image.open(headshot_path)
                    # Don't verify - just check if we can open it (verify is slow)
                    test_img.close()
                except Exception as e:
                    print(f"Warning: Headshot file is not a valid image: {e}, skipping headshot")
                    headshot_path = None
            
            headshot_paths = self._parse_headshots(headshot_path)
            headshot_paths = [p for p in headshot_paths if p and os.path.exists(p)]
            if not headshot_paths:
                print("   Skipping headshot processing (no valid headshot file)")
            else:
                # Remove background from headshot(s) to make them transparent
                use_api_removal = False
                try:
                    from config import Config
                    if hasattr(Config, 'REMOVEBG_API_KEY') and Config.REMOVEBG_API_KEY and Config.REMOVEBG_API_KEY.strip():
                        use_api_removal = True
                except:
                    pass
                
                print("   REMOVEBG_API_KEY not set, using manual background removal" if not use_api_removal else "   Using remove.bg API when available")

                # Headshot target box: wide, tall enough to reach near the slide bottom.
                headshot_area_width = 600
                headshot_area_height = 545
                
                # Position: directly below the map (just under Texas), same X as before.
                headshot_area_x = map_area_x + (map_width - headshot_area_width) // 2 - 30
                headshot_area_y = map_area_y + map_height + 10

                def load_process_headshot(path: str) -> Optional[Image.Image]:
                    try:
                        print(f"    Processing headshot: {path}")
                        
                        # 0) PREPARE ORIGINAL
                        original = Image.open(path).convert("RGBA")
                        # Keep original size reasonable for processing (780x780 target, so 1500 is plenty)
                        if max(original.size) > 1500:
                            original.thumbnail((1500, 1500), Image.Resampling.LANCZOS)
                        
                        # --- HELPER: MAP BACKGROUND FALLBACK ---
                        # If BG removal fails, composite person on top of map to simulate transparency
                        # Try to get map from slide if map_path is not available
                        def get_map_for_background():
                            """Try to get map image from various sources."""
                            # First, try bundled map background image (committed to repo)
                            script_dir = os.path.dirname(os.path.abspath(__file__))
                            bundled_map_paths = [
                                os.path.join(script_dir, "assets", "map_background.png"),
                                os.path.join(script_dir, "map_background.png"),
                                os.path.join(os.path.dirname(script_dir), "assets", "map_background.png"),
                            ]
                            for bundled_path in bundled_map_paths:
                                if os.path.exists(bundled_path):
                                    try:
                                        print(f"    Using bundled map background: {bundled_path}")
                                        return Image.open(bundled_path).convert("RGBA")
                                    except Exception as e:
                                        print(f"    Warning: Could not load bundled map: {e}")
                                        continue
                            
                            # Second, try map_path if available
                            if map_path and os.path.exists(map_path):
                                try:
                                    return Image.open(map_path).convert("RGBA")
                                except Exception as e:
                                    print(f"    Warning: Could not load map from map_path: {e}")
                            
                            # Third, try common temp locations
                            import tempfile
                            temp_dir = tempfile.gettempdir()
                            possible_map_paths = [
                                os.path.join(temp_dir, "map.png"),
                                os.path.join(temp_dir, "map_placeholder.png"),
                            ]
                            for possible_path in possible_map_paths:
                                if os.path.exists(possible_path):
                                    try:
                                        return Image.open(possible_path).convert("RGBA")
                                    except:
                                        continue
                            
                            return None
                        
                        def make_map_background(img_in):
                            print("    Applying Transparent Background Fallback...")
                            # 1. Create a square canvas based on smallest dimension
                            w, h = img_in.size
                            d = min(w, h)
                            # Center crop to square
                            left = (w - d) // 2
                            top = (h - d) // 2
                            img_sq = img_in.crop((left, top, left + d, top + d))
                            
                            # 2. Try aggressive background removal with very conservative settings
                            # Use the gray flood-fill with very low tolerance to only remove obvious background
                            try:
                                import numpy as np
                                img_rgba = img_sq.convert("RGBA")
                                arr = np.array(img_rgba)
                                H, W = arr.shape[:2]
                                
                                # Luminance
                                lum = (0.299 * arr[..., 0] + 0.587 * arr[..., 1] + 0.114 * arr[..., 2]).astype(np.float32)
                                alpha = arr[..., 3].astype(np.uint8)
                                
                                # Border samples for background color
                                bs = max(8, min(H, W) // 25)
                                border = np.concatenate([
                                    lum[:bs, :].ravel(),
                                    lum[-bs:, :].ravel(),
                                    lum[:, :bs].ravel(),
                                    lum[:, -bs:].ravel(),
                                ])
                                bg = np.median(border)
                                dist = np.abs(lum - bg)
                                
                                # Very conservative flood-fill - only remove pixels very close to border color
                                # Protect center 70% of image
                                center_y_start = int(H * 0.15)
                                center_y_end = int(H * 0.85)
                                center_x_start = int(W * 0.15)
                                center_x_end = int(W * 0.85)
                                center_protection = np.zeros((H, W), dtype=bool)
                                center_protection[center_y_start:center_y_end, center_x_start:center_x_end] = True
                                
                                from collections import deque
                                def flood(t):
                                    close = dist <= t
                                    bg_mask = np.zeros((H, W), dtype=bool)
                                    q = deque()
                                    
                                    def push(y, x):
                                        # Check bounds FIRST before accessing arrays
                                        if not (0 <= y < H and 0 <= x < W):
                                            return
                                        if center_protection[y, x]:
                                            return
                                        if close[y, x] and not bg_mask[y, x]:
                                            bg_mask[y, x] = True
                                            q.append((y, x))
                                    
                                    for x in range(W):
                                        push(0, x); push(H - 1, x)
                                    for y in range(H):
                                        push(y, 0); push(y, W - 1)
                                    
                                    while q:
                                        y, x = q.popleft()
                                        for dy in (-1, 0, 1):
                                            for dx in (-1, 0, 1):
                                                if dy == 0 and dx == 0:
                                                    continue
                                                push(y + dy, x + dx)
                                    
                                    return bg_mask
                                
                                # Try very low tolerance (3-8) to only remove obvious background
                                best_mask = None
                                for t in [3, 5, 8]:
                                    mask = flood(t)
                                    removed = mask.mean()
                                    if 0.05 <= removed <= 0.25:  # Only remove 5-25% (very conservative)
                                        best_mask = mask
                                        print(f"    ✓ Conservative background removal: tol={t}, removed={removed:.1%}")
                                        break
                                
                                if best_mask is None:
                                    # Use the most conservative mask
                                    best_mask = flood(3)
                                
                                # Apply mask to alpha channel
                                new_alpha = alpha.copy()
                                new_alpha[best_mask] = 0
                                
                                # Create result with transparent background
                                result = arr.copy()
                                result[..., 3] = new_alpha
                                result_img = Image.fromarray(result, "RGBA")
                                
                                # Convert to grayscale while preserving alpha
                                result_gray = result_img.convert("L")
                                alpha_channel = result_img.split()[3]
                                
                                return Image.merge("RGBA", (result_gray, result_gray, result_gray, alpha_channel))
                                
                            except Exception as e:
                                print(f"    ⚠️  Advanced removal failed: {e}, using circular crop with transparent background")
                                # Fallback: circular crop with transparent background
                                mask = Image.new('L', (d, d), 0)
                                draw_mask = ImageDraw.Draw(mask)
                                draw_mask.ellipse((0, 0, d, d), fill=255)
                                img_sq.putalpha(mask)
                                gray = img_sq.convert("L")
                                return Image.merge("RGBA", (gray, gray, gray, mask))

                        # 1) Attempt Background Removal
                        img = self._remove_bg_best_effort(path, use_api=use_api_removal)
                        
                        # Check if background removal returned None
                        if img is None:
                            print(f"    WARNING: Background removal returned None. Using Transparent Background Fallback.")
                            return make_map_background(original)

                        # 2) SAFETY CHECK 1: Did we wipe the image? 
                        # If opacity is < 10%, we deleted too much of the person. Use Circular Fallback.
                        opaque_frac, transp_frac, mean_alpha = self._alpha_stats(img)
                        print(f"    After BG removal: opaque={opaque_frac:.2f}, transp={transp_frac:.2f}, meanA={mean_alpha:.0f}")
                        if opaque_frac < 0.10: 
                            print(f"    WARNING: BG removal wiped image (opaque={opaque_frac:.2f}). Using Transparent Background Fallback.")
                            return make_map_background(original)
                        
                        # 2b) SAFETY CHECK: If we removed too much (<20% of image is person), that's suspicious
                        # Note: 40-50% opaque is actually good for headshots (person) with 50-60% transparent (background)
                        if opaque_frac < 0.20:
                            print(f"    WARNING: Only {opaque_frac:.1%} of image remains - may have removed too much. Using Transparent Background Fallback.")
                            return make_map_background(original)

                        # 3) Only clean up mask if rembg didn't work well (low transparency)
                        # If rembg already gave us good transparency, skip _fix_alpha_mask to avoid wiping it
                        if transp_frac < 0.10:
                            # Background removal was weak, try to clean it up
                            print(f"    Transparency is weak ({transp_frac:.2f}), attempting to clean mask...")
                            img_before_fix = img.copy()
                            img = self._fix_alpha_mask(img, a_min=2)
                            
                            # Check if _fix_alpha_mask made it worse
                            opaque_after, transp_after, mean_after = self._alpha_stats(img)
                            print(f"    After _fix_alpha_mask: opaque={opaque_after:.2f}, transp={transp_after:.2f}, meanA={mean_after:.0f}")
                            
                            # If it made it worse (less opaque), revert
                            if opaque_after < opaque_frac * 0.5:  # Lost more than 50% of opacity
                                print(f"    WARNING: _fix_alpha_mask made it worse, reverting...")
                                img = img_before_fix
                        else:
                            print(f"    Good transparency ({transp_frac:.2f}), skipping _fix_alpha_mask to preserve quality")
                        
                        # 4) Final safety check before processing
                        final_check_o, final_check_t, _ = self._alpha_stats(img)
                        if final_check_o < 0.05:
                            print(f"    WARNING: Image has no subject (opaque={final_check_o:.2f}). Using Transparent Background Fallback.")
                            return make_map_background(original)

                        # 5) Standard Processing (Grayscale + Edges)
                        img = self._to_grayscale_preserve_alpha(img)
                        img = self._refine_edges(img, erode_size=1, blur_radius=0.5)
                        
                        # 6) Harden alpha channel - ensure background is fully transparent (0) and foreground is fully opaque (255)
                        # This prevents semi-transparent halos when compositing
                        try:
                            import numpy as np
                            arr = np.array(img)
                            alpha = arr[..., 3].astype(np.uint8)
                            # Threshold: pixels with alpha < 128 become fully transparent (0), >= 128 become fully opaque (255)
                            # This creates a clean cutout without semi-transparent edges
                            alpha_hard = np.where(alpha < 128, 0, 255).astype(np.uint8)
                            arr[..., 3] = alpha_hard
                            img = Image.fromarray(arr, "RGBA")
                            print(f"    Hardened alpha channel (threshold=128) to ensure clean transparency")
                        except Exception as e:
                            print(f"    Warning: Could not harden alpha: {e}")
                        
                        # Final check before returning
                        final_o, final_t, final_ma = self._alpha_stats(img)
                        print(f"    Final headshot stats: opaque={final_o:.2f}, transp={final_t:.2f}, meanA={final_ma:.0f}")
                        if final_o < 0.05:
                            print(f"    WARNING: Final image has no subject (opaque={final_o:.2f}). Using Transparent Background Fallback.")
                            return make_map_background(original)

                        return img

                    except Exception as e:
                        print(f"Warning: failed headshot {path}: {e}")
                        import traceback
                        traceback.print_exc()
                        # Final safety: Return original with map background
                        try:
                            orig = Image.open(path).convert("RGBA")
                            return make_map_background(orig)
                        except:
                            return None

                imgs = [load_process_headshot(p) for p in headshot_paths[:2]]
                imgs = [im for im in imgs if im is not None]

                if len(imgs) == 1:
                    im = imgs[0]
                    # Ensure image has alpha channel for transparency
                    if im.mode != 'RGBA':
                        im = im.convert('RGBA')
                    im = self._resize_cover(im, headshot_area_width, headshot_area_height)
                    # Paste with alpha mask to preserve transparency
                    if im.mode == 'RGBA':
                        slide.paste(im, (headshot_area_x, headshot_area_y), im.split()[3])
                    else:
                        slide.paste(im, (headshot_area_x, headshot_area_y))
                    draw = ImageDraw.Draw(slide)
                elif len(imgs) == 2:
                    gap = 30
                    each_w = (headshot_area_width - gap) // 2
                    each_h = headshot_area_height
                    left, right = imgs
                    # Ensure both images have alpha channel for transparency
                    if left.mode != 'RGBA':
                        left = left.convert('RGBA')
                    if right.mode != 'RGBA':
                        right = right.convert('RGBA')
                    left = self._resize_cover(left, each_w, each_h)
                    right = self._resize_cover(right, each_w, each_h)

                    left_x = headshot_area_x
                    right_x = headshot_area_x + each_w + gap
                    y = headshot_area_y

                    # Paste with alpha mask to preserve transparency
                    if left.mode == 'RGBA':
                        slide.paste(left, (left_x, y), left.split()[3])
                    else:
                        slide.paste(left, (left_x, y))
                    if right.mode == 'RGBA':
                        slide.paste(right, (right_x, y), right.split()[3])
                    else:
                        slide.paste(right, (right_x, y))
                draw = ImageDraw.Draw(slide)
        except Exception as e:
            print(f"Warning: Could not load headshot: {e}")
            import traceback
            traceback.print_exc()
        
        # 8. Replace investment stage in sidebar (transparent, black, thicker text like company name, includes year)
        # Find Slauson&Co position in sidebar (bottom of sidebar)
        slauson_y = height - 200  # Approximate position of "SLAUSON&CO." text (bottom of sidebar)
        
        # Parse investment stage to match the template look (e.g., "SEED Q4, 2024")
        # Input format is typically like "SEED Q2, 2024" or "PRE-SEED Q2, 2024"
        stage_parts = investment_stage.upper().split(',')
        if len(stage_parts) >= 2:
            stage_quarter = stage_parts[0].strip()  # "SEED Q2" or "PRE-SEED Q2"
            year_text = stage_parts[1].strip()  # "2024"
            # Combine into "STAGE Q#, YYYY" format
            stage_text = f"{stage_quarter}, {year_text}"  # "SEED Q2, 2024"
        else:
            stage_text = investment_stage.upper()
        
        # --- Render investment stage (fix "squished" look) ---
        # Keep the old font behavior (no preferred_family), but avoid heavy multi-pass stroke
        # and avoid large downscales after rotation (both make the text look condensed).
        sidebar_w = 200
        # Keep the sidebar stage text readable but not comically large.
        # Width (after rotation) is driven mostly by font height, so side margins + padding
        # are the safest knobs to make it "reasonable" without reintroducing squish/blur.
        top_margin = 35
        bottom_margin = 25
        # Constrain usable width so the font size matches the template (like your screenshot).
        side_margin = 35
        max_w = sidebar_w - 2 * side_margin
        max_h = height - top_margin - bottom_margin

        stage_color = (0, 0, 0)
        # Add back some padding so the text doesn't dominate the sidebar.
        padding = 10

        stage_img = None
        final_w = final_h = 0

        # Choose the largest font size that fits without resizing the rotated bitmap.
        # Cap the max size so it stays visually consistent with the template.
        for fs in range(55, 11, -1):
            # Stage text should use the same font family as the title (Bebas Neue).
            stage_font = self._load_font(fs, bold=False, preferred_family="BebasNeue")

            tmp = Image.new("RGBA", (2000, 2000), (0, 0, 0, 0))
            td = ImageDraw.Draw(tmp)
            bbox = td.textbbox((0, 0), stage_text, font=stage_font)
            text_w = bbox[2] - bbox[0]
            text_h = bbox[3] - bbox[1]

            img = Image.new("RGBA", (max(1, int(text_w + padding * 2)), max(1, int(text_h + padding * 2))), (0, 0, 0, 0))
            d = ImageDraw.Draw(img)
            d.text(
                (padding - bbox[0], padding - bbox[1]),
                stage_text,
                fill=stage_color,
                font=stage_font,
                stroke_width=1,
                stroke_fill=stage_color,
            )

            rot = img.rotate(90, expand=True, resample=Image.Resampling.BICUBIC)
            w, h = rot.size
            if w <= max_w and h <= max_h:
                stage_img = rot
                final_w, final_h = w, h
                break

        if stage_img is None:
            # Last resort: render at a small size; still no post-rotate downscale.
            stage_font = self._load_font(18, bold=False)
            tmp = Image.new("RGBA", (2000, 2000), (0, 0, 0, 0))
            td = ImageDraw.Draw(tmp)
            bbox = td.textbbox((0, 0), stage_text, font=stage_font)
            text_w = bbox[2] - bbox[0]
            text_h = bbox[3] - bbox[1]
            img = Image.new("RGBA", (max(1, int(text_w + padding * 2)), max(1, int(text_h + padding * 2))), (0, 0, 0, 0))
            d = ImageDraw.Draw(img)
            d.text(
                (padding - bbox[0], padding - bbox[1]),
                stage_text,
                fill=stage_color,
                font=stage_font,
                stroke_width=1,
                stroke_fill=stage_color,
            )
            stage_img = img.rotate(90, expand=True, resample=Image.Resampling.BICUBIC)
            final_w, final_h = stage_img.size

        # Compute placement
        # Center the stage text in the sidebar and place it near the top like the template.
        paste_x = max(0, (sidebar_w - final_w) // 2)
        paste_y = 110

        # Clamp so it can NEVER go off-canvas
        paste_x = max(0, min(paste_x, sidebar_w - final_w)) - 10
        paste_y = max(0, min(paste_y, height - final_h)) - 10

        # Clear the existing baked-in sidebar stage text under our overlay.
        # The template already contains stage text; drawing on top can look "squished"/blurry.
        try:
            # Sample the sidebar background (exclude black text) to get the orange fill.
            sidebar_bg_rgb = self._get_dominant_color(
                template,
                region=(10, 300, 180, 300),
                exclude_colors=[(0, 0, 0), (255, 255, 255)],
            )
            x0 = max(0, paste_x - 4)
            y0 = max(0, paste_y - 4)
            x1 = min(sidebar_w, paste_x + final_w + 4)
            y1 = min(height, paste_y + final_h + 4)
            draw.rectangle([(x0, y0), (x1, y1)], fill=sidebar_bg_rgb)
        except Exception:
            pass

        slide.paste(stage_img, (paste_x, paste_y), stage_img)

        # --- Per-slide tracking watermark (visible, for human traceability) ---
        # Note: The PDF output from this generator is image-based, so this watermark is not machine-extractable
        # via PDF text extraction. We use a separate index store for reliable replacement/deletion.
        slide_job_id = (
            company_data.get("slide_job_id")
            or company_data.get("Slide Job ID")
            or company_data.get("job_id")
        )
        if slide_job_id:
            try:
                watermark_draw = ImageDraw.Draw(slide)
                wm_font = self._load_font(14, bold=False)
                wm_text = f"Slide Job ID: {slide_job_id}"
                # Bottom-right placement (very right)
                wm_bbox = watermark_draw.textbbox((0, 0), wm_text, font=wm_font)
                wm_w = wm_bbox[2] - wm_bbox[0]
                wm_h = wm_bbox[3] - wm_bbox[1]
                margin = 10
                wm_x = max(margin, width - wm_w - margin)
                wm_y = max(margin, height - wm_h - margin)

                # Pure solid black
                wm_fill = (0, 0, 0, 255) if slide.mode == "RGBA" else (0, 0, 0)
                watermark_draw.text((wm_x, wm_y), wm_text, fill=wm_fill, font=wm_font)
            except Exception:
                pass
        
        # Convert to RGB for PDF
        slide_rgb = Image.new('RGB', slide.size, (42, 42, 42))
        slide_rgb.paste(slide, mask=slide.split()[3] if slide.mode == 'RGBA' else None)
        
        # Return based on output format
        if output_format.lower() == "pptx":
            return self._create_pptx_from_slide(
                slide, company_data, headshot_path, logo_path, map_path,
                map_area_x, map_area_y, map_width, map_height
            )
        
        # Default: Convert to PDF (flattened, not editable)
        # Convert to RGB for PDF
        slide_rgb = Image.new('RGB', slide.size, (42, 42, 42))
        slide_rgb.paste(slide, mask=slide.split()[3] if slide.mode == 'RGBA' else None)
        
        # Convert to PDF
        # Use temporary file to avoid PIL's fileno() issue with BytesIO
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
            try:
                slide_rgb.save(tmp_file.name, format='PNG', optimize=False)
                with open(tmp_file.name, 'rb') as f:
                    img_bytes = f.read()
                pdf_bytes = img2pdf.convert(img_bytes)
                return pdf_bytes
            except Exception as e:
                print(f"Error: {e}")
                return None
            finally:
                # Clean up temporary file
                try:
                    os.unlink(tmp_file.name)
                except:
                    pass
    
    def _create_pptx_from_slide(
        self, 
        template_slide: Image.Image,
        company_data: Dict,
        headshot_path: str,
        logo_path: str,
        map_path: Optional[str],
        map_area_x: int,
        map_area_y: int,
        map_width: int,
        map_height: int
    ) -> bytes:
        """
        Create an editable PPTX file with separate elements for Canva import.
        This allows headshot, logo, map, and text to be editable in Canva.
        
        Args:
            template_slide: PIL Image of the template background
            company_data: Company information
            headshot_path: Path to headshot image
            logo_path: Path to logo image
            map_path: Path to map image (optional)
            map_area_x, map_area_y, map_width, map_height: Map position/size
            
        Returns:
            PPTX file bytes
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is required for PPTX generation. Install with: pip install python-pptx")
        
        # Slide dimensions: 1920x1080 pixels = 20x11.25 inches at 96 DPI
        # Standard widescreen: 13.333" x 7.5" (16:9)
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 1920px / 144 DPI
        prs.slide_height = Inches(7.5)   # 1080px / 144 DPI
        
        # Create blank slide
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide_pptx = prs.slides.add_slide(slide_layout)
        
        # CRITICAL: Don't add full-slide background image - it causes Canva to flatten everything into one image
        # Instead, we'll add ONLY the editable elements (logo, headshot, text, map) on a blank slide
        # The user can add their own background in Canva, or we can add it as a separate editable element later
        # For now, just add a simple colored background using shapes (editable)
        from pptx.enum.shapes import MSO_SHAPE
        
        # 1) Orange sidebar (left side, editable shape)
        sidebar_width_inch = Inches(200 / 96.0)  # 200px converted to inches
        sidebar = slide_pptx.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            sidebar_width_inch,
            prs.slide_height
        )
        sidebar.fill.solid()
        sidebar.fill.fore_color.rgb = RGBColor(242, 140, 40)  # Orange #F28C28
        sidebar.line.fill.background()  # No border
        
        # 2) Dark grey main area (right side, editable shape)
        main_area = slide_pptx.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            sidebar_width_inch,
            Inches(0),
            prs.slide_width - sidebar_width_inch,
            prs.slide_height
        )
        main_area.fill.solid()
        main_area.fill.fore_color.rgb = RGBColor(42, 42, 42)  # Dark grey
        main_area.line.fill.background()  # No border
        
        # Get actual positions from template (convert pixels to inches at 96 DPI)
        width, height = template_slide.size
        px_to_inch = 1.0 / 96.0  # 96 DPI standard
        
        # Add "SLAUSON&CO." text in sidebar (bottom, editable text)
        slauson_text = "SLAUSON&CO."
        slauson_y_px = height - 200  # Bottom of sidebar
        slauson_x_px = 10  # Left side of sidebar
        
        slauson_textbox = slide_pptx.shapes.add_textbox(
            Inches(slauson_x_px * px_to_inch),
            Inches(slauson_y_px * px_to_inch),
            Inches(180 * px_to_inch),  # Width for sidebar
            Inches(1.5)
        )
        slauson_tf = slauson_textbox.text_frame
        slauson_tf.word_wrap = False
        slauson_p = slauson_tf.paragraphs[0]
        slauson_run = slauson_p.add_run()
        slauson_run.text = slauson_text
        slauson_run.font.size = Pt(28)
        slauson_run.font.bold = True
        slauson_run.font.color.rgb = RGBColor(0, 0, 0)  # Black
        
        # 2) Logo (top right, editable) - from PIL code: logo_x, logo_y = width - 170, 10
        if logo_path and os.path.exists(logo_path):
            logo_img = Image.open(logo_path).convert('RGBA')
            # Load image fully before operations to avoid lazy loading issues
            logo_img.load()
            # Make circular (same as PIL code)
            logo_size_px = 130
            
            # Create circular mask first
            circle_mask = Image.new('L', (logo_size_px, logo_size_px), 0)
            circle_draw = ImageDraw.Draw(circle_mask)
            circle_draw.ellipse([(0, 0), (logo_size_px, logo_size_px)], fill=255)
            
            # Center crop to square first
            logo_w, logo_h = logo_img.size
            min_dim = min(logo_w, logo_h)
            logo_img = logo_img.crop(((logo_w - min_dim) // 2, (logo_h - min_dim) // 2, 
                                     (logo_w + min_dim) // 2, (logo_h + min_dim) // 2))
            
            # Resize to fill the entire circle (logo_size_px x logo_size_px)
            logo_img = logo_img.resize((logo_size_px, logo_size_px), Image.Resampling.LANCZOS)
            
            # Apply circular mask to logo - this makes it truly circular
            logo_masked = Image.new('RGBA', (logo_size_px, logo_size_px), (0, 0, 0, 0))
            logo_masked.paste(logo_img, (0, 0))  # Paste at origin, no padding
            logo_masked.putalpha(circle_mask)  # Apply circular mask to make it round
            
            logo_bytes = io.BytesIO()
            logo_masked.save(logo_bytes, format='PNG')
            logo_bytes.seek(0)
            
            # Position: from PIL code - logo_x, logo_y = width - 170, 10
            logo_x_inch = (width - 170) * px_to_inch
            logo_y_inch = 10 * px_to_inch
            logo_size_inch = logo_size_px * px_to_inch
            
            slide_pptx.shapes.add_picture(
                logo_bytes,
                Inches(logo_x_inch),
                Inches(logo_y_inch),
                width=Inches(logo_size_inch)
            )
        
        # 3) Company name (editable text) - from PIL code: name_x, name_y = founders_text_x - 50, 120
        company_name = company_data.get("name", "").upper()
        if company_name:
            # Position from PIL: name_x = founders_text_x - 50, name_y = 120
            # founders_text_x = 320 from PIL code
            name_x_px = 320 - 50  # 270
            name_y_px = 120
            
            # Calculate font size (same logic as PIL)
            base_font_size = 200
            name_font_size = base_font_size
            # Check for overlap with map
            max_allowed_width = map_area_x - name_x_px - 50
            # Estimate text width (rough calculation)
            estimated_width = len(company_name) * (name_font_size * 0.6)  # Rough estimate
            if estimated_width > max_allowed_width:
                name_font_size = int((max_allowed_width / estimated_width) * base_font_size)
                # Prefer a larger minimum, but only if it still fits by estimate.
                desired_min_size = 120
                if name_font_size < desired_min_size:
                    desired_estimated_width = len(company_name) * (desired_min_size * 0.6)
                    if desired_estimated_width <= max_allowed_width:
                        name_font_size = desired_min_size
                name_font_size = max(90, name_font_size)
            
            textbox = slide_pptx.shapes.add_textbox(
                Inches(name_x_px * px_to_inch),
                Inches(name_y_px * px_to_inch),
                Inches((map_area_x - name_x_px - 50) * px_to_inch),
                Inches(2.0)  # Height for large text
            )
            tf = textbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = company_name
            run.font.size = Pt(name_font_size)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 140, 0)  # Orange color
        
        # 4) Investment stage (editable text) - positioned in sidebar
        investment_stage = company_data.get("investment_stage", "")
        if not investment_stage:
            round_val = company_data.get("investment_round", "PRE-SEED")
            quarter_val = company_data.get("quarter", "Q2")
            year_val = company_data.get("year", "2024")
            investment_stage = f"{round_val} {quarter_val} {year_val}"
        
        if investment_stage:
            paste_x_px = 15   # slightly right
            paste_y_px = 70   # lower
            
            # Make the textbox TALL/WIDE BEFORE rotation so it has room after rotation.
            # Use most of the slide height.
            box_w_px = 180
            box_h_px = 950   # <-- big, prevents "tiny" rendering
            
            textbox = slide_pptx.shapes.add_textbox(
                Inches(paste_x_px * px_to_inch),
                Inches(paste_y_px * px_to_inch),
                Inches(box_w_px * px_to_inch),
                Inches(box_h_px * px_to_inch),
            )
            textbox.rotation = 270  # bottom->top like the sidebar
            
            tf = textbox.text_frame
            tf.word_wrap = False
            
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = investment_stage
            
            # BIGGER FONT (36-46pt range)
            run.font.size = Pt(40)
            run.font.bold = False  # Less bold for easier reading
            run.font.color.rgb = RGBColor(0, 0, 0)  # Black
        
        # 5) Headshot (below map, editable)
        if headshot_path and os.path.exists(headshot_path):
            headshot_bytes = io.BytesIO()
            headshot_img = Image.open(headshot_path).convert('RGBA')
            headshot_img.load()  # Load fully before save to avoid memory issues
            # Ensure transparent background for PPTX as well (grayscale-aware)
            try:
                headshot_img = self._remove_background_gray(headshot_img, tol=8, feather=1)
            except Exception as e:
                print(f"   Warning: PPTX fallback background removal failed: {e}")
            # ALWAYS convert to grayscale while preserving alpha (headshot should be grayscale)
            headshot_img = self._to_grayscale_preserve_alpha(headshot_img)
            # Process headshot (background removal already done in PIL code)
            headshot_img.save(headshot_bytes, format='PNG', optimize=False)
            headshot_bytes.seek(0)
            
            # Position from PIL: headshot_area_width = 550 * 2.2 = 1210, headshot_area_height = 500 * 2.2 = 1100
            # headshot_area_x = map_area_x + (map_width - headshot_area_width) // 2 - 50
            # headshot_area_y = map_area_y + map_height - 50
            # Headshot target box (smaller to not cover map)
            # Keep PPTX placement consistent with the PDF layout (just under Texas)
            headshot_area_width_px = 600
            headshot_area_height_px = 545
            headshot_area_x_px = map_area_x + (map_width - headshot_area_width_px) // 2 - 30
            headshot_area_y_px = map_area_y + map_height + 10
            
            slide_pptx.shapes.add_picture(
                headshot_bytes,
                Inches(headshot_area_x_px * px_to_inch),
                Inches(headshot_area_y_px * px_to_inch),
                width=Inches(headshot_area_width_px * px_to_inch),
                height=Inches(headshot_area_height_px * px_to_inch)
            )
        
        # 6) Founders text (editable) - from PIL code: founders_text_x = 320, founders_text_y = 415
        founders_items = self._parse_name_list(company_data.get("founders", ""))[:3]
        founders_text = "\n".join(founders_items)
        if founders_text:
            
            # Position from PIL: founders_text_x = 320, founders_text_y = 415
            founders_text_x_px = 320
            founders_text_y_px = 415
            
            textbox = slide_pptx.shapes.add_textbox(
                Inches(founders_text_x_px * px_to_inch),
                Inches(founders_text_y_px * px_to_inch),
                Inches(3.0),
                Inches(2.0)
            )
            tf = textbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = founders_text
            run.font.size = Pt(24)
            run.font.color.rgb = RGBColor(255, 255, 255)  # White
        
        # 7) Co-investors text (editable) - from PIL code: investors_text_x = 650, investors_text_y = 415
        co_investors_items = self._parse_name_list(company_data.get("co_investors", ""))[:4]
        co_investors_text = "\n".join(co_investors_items)
        if co_investors_text:
            
            # Position from PIL: investors_text_x = 650, investors_text_y = 415
            investors_text_x_px = 650
            investors_text_y_px = 415
            
            textbox = slide_pptx.shapes.add_textbox(
                Inches(investors_text_x_px * px_to_inch),
                Inches(investors_text_y_px * px_to_inch),
                Inches(3.0),
                Inches(2.0)
            )
            tf = textbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = co_investors_text
            run.font.size = Pt(24)
            run.font.color.rgb = RGBColor(255, 255, 255)  # White
        
        # 8) Background text (editable) - from PIL code: bg_text_x = 320, bg_text_y = 650
        background_text = company_data.get("background", company_data.get("description", ""))
        if background_text:
            # Position from PIL: bg_text_x = 320, bg_text_y = 650, bg_text_width = 700
            bg_text_x_px = 320
            bg_text_y_px = 650
            bg_text_width_px = 700
            
            textbox = slide_pptx.shapes.add_textbox(
                Inches(bg_text_x_px * px_to_inch),
                Inches(bg_text_y_px * px_to_inch),
                Inches(bg_text_width_px * px_to_inch),
                Inches(2.5)
            )
            tf = textbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = background_text
            run.font.size = Pt(32)
            run.font.color.rgb = RGBColor(255, 255, 255)  # White
        
        # 9) Map (if provided, editable) - positioned at detected map area
        if map_path and os.path.exists(map_path):
            map_bytes = io.BytesIO()
            map_img = Image.open(map_path).convert('RGBA')
            map_img.load()  # Load fully before save to avoid memory issues
            map_img.save(map_bytes, format='PNG', optimize=False)
            map_bytes.seek(0)
            
            # Position: top right area (from detected map area)
            map_x = map_area_x / 96.0  # Convert pixels to inches
            map_y = map_area_y / 96.0
            map_w = map_width / 96.0
            map_h = map_height / 96.0
            
            slide_pptx.shapes.add_picture(
                map_bytes,
                Inches(map_x),
                Inches(map_y),
                width=Inches(map_w),
                height=Inches(map_h)
            )
        
        # Save to bytes
        pptx_bytes = io.BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        return pptx_bytes.getvalue()
